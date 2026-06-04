"""
pymims_app.py  —  PyMIMS v1.0 desktop GUI (Panel)
=================================================================

    Imaging   — single channel / ratio / HSI; display filters, a GLOBAL
                binning (sum-pool) control applied uniformly to every image so
                overlays/side-by-sides co-register, HSI ratio-median, Plane QC
                (per-plane scan + auto/manual drop), and save-at-DPI
    HMR       — high mass resolution deconvolution            (placeholder)
    Analysis  — clustering + intensity histograms             (placeholder)
    Draw ROIs — interactive hand-drawn ROIs → Analysis stats   (Bokeh canvas)
    Metadata  — full reverse-engineered Cameca header         (working)

Run:  panel serve pymims_app.py --show
(--autoreload triggers a Panel warm-up bug on some versions; omit it.)
"""

import os
import io
import sys
import contextlib
import datetime

# Ensure this app's own directory (the repo root) is importable, so the
# top-level HMR engine modules (irf/deconvolve/sparse_deconvolve/hmr_identify)
# resolve even when `pymims` itself is pip-installed and the repo root is not
# otherwise on sys.path (e.g. under `panel serve`).
for _p in (
    os.path.dirname(os.path.abspath(__file__)) if "__file__" in globals() else None,
    os.getcwd(),
):
    if _p and _p not in sys.path:
        sys.path.insert(0, _p)

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np

import panel as pn
import param

from pymims import MimsImage, save_figure
from pymims import metadata as pymims_metadata
from pymims.core import _bin_sum            # same crop+sum-pool helper as the library

# Treat the GUI as an inline-display host: plot_* methods return their Figure
# without writing PNGs to disk or closing them.
import pymims.core as _pymims_core
_pymims_core._IS_NOTEBOOK = True

# Analysis-tab backends. Canonical layout is the package submodule form
# (matching each module's own `from .histograms import ...` relative import);
# the flat fallback covers a non-package checkout.
try:
    from pymims.clustering import (
        cluster_pixels, plot_cluster_labels, plot_cluster_grid,
        plot_overlay, plot_metric_sweep, plot_dendrogram,
        extract_cluster_masks,
    )
    from pymims.histograms import plot_histograms
except ImportError:                                   # flat-layout fallback
    from pymims_clustering import (
        cluster_pixels, plot_cluster_labels, plot_cluster_grid,
        plot_overlay, plot_metric_sweep, plot_dendrogram,
        extract_cluster_masks,
    )
    from pymims_histograms import plot_histograms

# HMR engine (irf + deconvolve + sparse_deconvolve + hmr_identify). These are
# standalone modules dropped in the repo root; import is best-effort so the
# rest of the app still runs if they're absent.
try:
    from irf import rect_gauss, fit_sigma_from_edge
    from deconvolve import deconvolve, deconvolve_with_N_sweep
    from sparse_deconvolve import deconvolve_sparse
    from hmr_identify import (enumerate_candidates, match_by_separation,
                              identify_by_difference, identify_blind,
                              ELEMENTS as HMR_ELEMENTS)
    _HMR_OK = True
    _HMR_ERR = ""
except Exception as _hmr_exc:                              # pragma: no cover
    try:
        from pymims.irf import rect_gauss, fit_sigma_from_edge
        from pymims.deconvolve import deconvolve, deconvolve_with_N_sweep
        from pymims.sparse_deconvolve import deconvolve_sparse
        from pymims.hmr_identify import (enumerate_candidates, match_by_separation,
                                         identify_by_difference, identify_blind,
                                         ELEMENTS as HMR_ELEMENTS)
        _HMR_OK = True
        _HMR_ERR = ""
    except Exception as _hmr_exc2:
        _HMR_OK = False
        # report the top-level (flat) failure — that's the informative one
        _HMR_ERR = f"{type(_hmr_exc).__name__}: {_hmr_exc}"

try:
    from pymims import __version__ as PYMIMS_VERSION
except Exception:
    PYMIMS_VERSION = "1.0.0"

pn.extension("plotly", sizing_mode="stretch_width", notifications=True)

_BIN_OPTIONS = {"1× (full res)": 1, "2×2": 2, "4×4": 4, "8×8": 8}


# ─────────────────────────────────────────────────────────────────────────────
# Shared application state
# ─────────────────────────────────────────────────────────────────────────────

class State(param.Parameterized):
    img = param.Parameter(default=None)
    drift_version = param.Integer(default=0)   # bump to re-render after drift
    plane_version = param.Integer(default=0)   # bump to re-render after a plane drop
    bin_n = param.Integer(default=1)           # GLOBAL binning (sum-pool) factor


state = State()


# ─────────────────────────────────────────────────────────────────────────────
# Sidebar: file picker + global binning (sum-pool) control
# ─────────────────────────────────────────────────────────────────────────────

dir_input = pn.widgets.TextInput(name="Directory", value=os.path.expanduser("~"))
up_btn = pn.widgets.Button(name="⬆ Up one level", button_type="default", width=140)
subdir_select = pn.widgets.Select(name="Enter subfolder", options=["—"], value="—")
file_select = pn.widgets.Select(name=".im files here", options=[])
load_btn = pn.widgets.Button(name="Load file", button_type="primary", icon="upload")
load_status = pn.pane.Markdown("*No file loaded.*")

bin_select = pn.widgets.Select(name="Binning (sum-pool)",
                               options=list(_BIN_OPTIONS), value="1× (full res)")
bin_note = pn.pane.Markdown(
    "<span style='font-size:11px;opacity:0.75'>Sum-pools counts into n×n blocks "
    "and applies to **all** images (single channel, ratio, HSI) so they "
    "co-register. Non-destructive — the full-resolution data is never "
    "altered.</span>")


def _on_bin(event):
    state.bin_n = _BIN_OPTIONS.get(event.new, 1)
bin_select.param.watch(_on_bin, "value")


# Session save / load (.npz) — resume analysis without the original .im
session_dir = pn.widgets.TextInput(name="Session folder",
                                    value=os.path.expanduser("~"))
session_name = pn.widgets.TextInput(name="Save as (name)", value="session")
save_session_btn = pn.widgets.Button(name="💾 Save session",
                                      button_type="default")
session_select = pn.widgets.Select(name="Open session", options=[])
load_session_btn = pn.widgets.Button(name="📂 Load session",
                                      button_type="default")
session_status = pn.pane.Markdown("")
session_note = pn.pane.Markdown(
    "<span style='font-size:11px;opacity:0.75'>Saves the cleaned / drift-"
    "corrected stack, dropped-plane record, metadata and raw header to a "
    "<code>.npz</code> — reload it later without the original <code>.im</code>. "
    "“Open session” lists the <code>.npz</code> files in the session folder; "
    "pick one and Load. Binning is a view setting and is not saved.</span>")


def _session_folder():
    return os.path.expanduser(session_dir.value or "~")


def _refresh_sessions(*events):
    folder = _session_folder()
    try:
        npzs = sorted(f for f in os.listdir(folder)
                      if f.lower().endswith(".npz"))
    except Exception:
        npzs = []
    session_select.options = npzs
    if session_select.value not in npzs:
        session_select.value = npzs[0] if npzs else None


def _save_session(event):
    if state.img is None:
        session_status.object = "⚠️ Load a file before saving a session."
        return
    try:
        folder = _session_folder()
        os.makedirs(folder, exist_ok=True)
        name = (session_name.value or "session").strip()
        if not name.lower().endswith(".npz"):
            name += ".npz"
        path = os.path.join(folder, name)
        state.img.save_session(path, verbose=False)
        mb = os.path.getsize(path) / 1e6
        _refresh_sessions()
        session_select.value = name                # surface the just-saved file
        session_status.object = f"✅ Saved `{name}` ({mb:.1f} MB)"
    except Exception as exc:
        session_status.object = f"❌ Save failed: `{exc}`"


def _load_session(event):
    name = session_select.value
    if not name:
        session_status.object = ("⚠️ No session selected — pick one from "
                                 "“Open session”.")
        return
    try:
        path = os.path.join(_session_folder(), name)
        if not os.path.exists(path):
            session_status.object = f"⚠️ Not found: `{path}`"
            _refresh_sessions()
            return
        loaded = MimsImage.load_session(path, verbose=False)
        state.img = loaded                       # rebuilds all tabs
        m = loaded.metadata
        drift = "drift applied" if loaded.corrected is not None else "no drift"
        session_status.object = (
            f"✅ Loaded `{name}`  \n"
            f"{m.get('width', '?')}×{m.get('height', '?')} px · "
            f"{m.get('n_planes', '?')} planes · {drift}")
        load_status.object = (
            f"📂 **Session: {name}**  \n"
            f"(from {os.path.basename(loaded.path)})")
    except Exception as exc:
        session_status.object = f"❌ Load failed: `{exc}`"


session_dir.param.watch(_refresh_sessions, "value")
save_session_btn.on_click(_save_session)
load_session_btn.on_click(_load_session)
_refresh_sessions()


# ─────────────────────────────────────────────────────────────────────────────
# Journal: collect rendered views in memory, export to Word / PowerPoint
# ─────────────────────────────────────────────────────────────────────────────

JOURNAL = []   # list of entries, each {'caption': str, 'added': str} plus EITHER
               #   'png' : bytes   (a figure entry), OR
               #   'text': str     (a text/summary entry)


def _journal_add_text(text, caption):
    """Append a text/summary entry to the journal (no figure)."""
    JOURNAL.append({
        "text": text,
        "caption": caption,
        "added": datetime.datetime.now().strftime("%Y-%m-%d %H:%M"),
    })
    _update_journal_status()

journal_status = pn.pane.Markdown("Journal is empty.")
journal_dir = pn.widgets.TextInput(name="Export folder",
                                   value=os.path.expanduser("~"))
journal_name = pn.widgets.TextInput(name="Export name", value="pymims_journal")
journal_fmt = pn.widgets.Select(
    name="Format", value="PowerPoint (.pptx)",
    options=["PowerPoint (.pptx)", "Word (.docx)", "Both"])
export_journal_btn = pn.widgets.Button(name="📄 Export journal",
                                       button_type="primary")
remove_last_btn = pn.widgets.Button(name="Remove last", button_type="default")
clear_journal_btn = pn.widgets.Button(name="Clear journal", button_type="default")
export_status = pn.pane.Markdown("")
journal_note = pn.pane.Markdown(
    "<span style='font-size:11px;opacity:0.75'>Use the “➕ Add … to journal” "
    "buttons on the Imaging and Analysis tabs to collect figures (with auto "
    "captions) and text summaries, then export them here. PowerPoint lays "
    "figures out 6 per slide (2×3) and gives each text summary its own slide; "
    "Word stacks one item per block. Each export is stamped with the "
    "date.</span>")


def _update_journal_status():
    if not JOURNAL:
        journal_status.object = "Journal is empty."
        return
    lines = [f"**Journal — {len(JOURNAL)} item(s):**"]
    for i, e in enumerate(JOURNAL, 1):
        cap = e['caption']
        cap = cap if len(cap) <= 64 else cap[:61] + "…"
        tag = "📝 " if "text" in e else ""
        lines.append(f"{i}. {tag}{cap}")
    journal_status.object = "  \n".join(lines)


def _build_docx(path):
    from docx import Document
    from docx.shared import Inches, Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()
    doc.add_heading("PyMIMS Journal", level=0)
    stamp = datetime.datetime.now().strftime("%Y-%m-%d %H:%M")
    sub = doc.add_paragraph(f"Generated {stamp}  ·  {len(JOURNAL)} items")
    if sub.runs:
        sub.runs[0].italic = True

    for i, e in enumerate(JOURNAL, 1):
        if "text" in e:
            head = doc.add_paragraph(f"Note {i}.  {e['caption']}")
            for run in head.runs:
                run.bold = True
                run.font.size = Pt(10)
            body = doc.add_paragraph(e["text"])
            for run in body.runs:
                run.font.name = "Consolas"
                run.font.size = Pt(8.5)
        else:
            doc.add_picture(io.BytesIO(e['png']), width=Inches(6.0))
            doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
            cap = doc.add_paragraph(f"Figure {i}.  {e['caption']}")
            cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for run in cap.runs:
                run.italic = True
                run.font.size = Pt(9)
        doc.add_paragraph("")     # spacer between items
    doc.save(path)
    return path


def _build_pptx(path):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)      # 16:9
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    SW, SH = int(prs.slide_width), int(prs.slide_height)
    margin = int(Inches(0.3))
    title_h = int(Inches(0.55))
    gap = int(Inches(0.12))
    cap_h = int(Inches(0.38))
    cols, rows = 3, 2
    per = cols * rows

    grid_w = SW - 2 * margin
    grid_h = SH - title_h - 2 * margin
    cell_w = (grid_w - (cols - 1) * gap) // cols
    cell_h = (grid_h - (rows - 1) * gap) // rows
    img_h = cell_h - cap_h

    stamp = datetime.datetime.now().strftime("%Y-%m-%d %H:%M")

    # Build an ordered page list: consecutive figure entries are batched into
    # 6-up grid pages; each text entry becomes its own full-width text page.
    # This preserves the order in which items were added to the journal.
    pages = []   # ('figs', [(idx, entry), ...]) | ('text', (idx, entry))
    fig_batch = []
    for idx, e in enumerate(JOURNAL):
        if "text" in e:
            if fig_batch:
                pages.append(("figs", fig_batch)); fig_batch = []
            pages.append(("text", (idx, e)))
        else:
            fig_batch.append((idx, e))
            if len(fig_batch) == per:
                pages.append(("figs", fig_batch)); fig_batch = []
    if fig_batch:
        pages.append(("figs", fig_batch))

    n_pages = max(len(pages), 1)

    def _add_title(slide, page_no):
        tb = slide.shapes.add_textbox(margin, int(Inches(0.08)), grid_w, title_h)
        tf = tb.text_frame
        tf.text = f"PyMIMS Journal — {stamp}   (slide {page_no} of {n_pages})"
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True

    for page_no, (kind, payload) in enumerate(pages, 1):
        slide = prs.slides.add_slide(blank)
        _add_title(slide, page_no)

        if kind == "text":
            idx, e = payload
            box = slide.shapes.add_textbox(
                margin, title_h + margin, grid_w, grid_h)
            tf = box.text_frame
            tf.word_wrap = True
            head = tf.paragraphs[0]
            head.text = f"{idx + 1}. {e['caption']}"
            head.font.size = Pt(13)
            head.font.bold = True
            body = tf.add_paragraph()
            body.text = e["text"]
            # Monospace so threshold tables and centroid columns line up.
            for run in body.runs:
                run.font.name = "Consolas"
                run.font.size = Pt(10)
            continue

        # Figure grid page
        for k, (idx, e) in enumerate(payload):
            r, c = divmod(k, cols)
            left = margin + c * (cell_w + gap)
            top = title_h + margin + r * (cell_h + gap)

            pic = slide.shapes.add_picture(io.BytesIO(e['png']), left, top,
                                           width=cell_w)
            if pic.height > img_h:
                pic.width = int(pic.width * img_h / pic.height)
                pic.height = int(img_h)
            pic.left = int(left + (cell_w - pic.width) / 2)   # centre in cell

            cb = slide.shapes.add_textbox(left, top + img_h, cell_w, cap_h)
            ctf = cb.text_frame
            ctf.word_wrap = True
            ctf.text = f"{idx + 1}. {e['caption']}"
            p = ctf.paragraphs[0]
            p.font.size = Pt(7)
            p.alignment = PP_ALIGN.CENTER

    prs.save(path)
    return path


def _export_journal(event):
    if not JOURNAL:
        export_status.object = "⚠️ Journal is empty — add a view first."
        return
    folder = os.path.expanduser(journal_dir.value or "~")
    base = (journal_name.value or "pymims_journal").strip()
    want_docx = journal_fmt.value in ("Word (.docx)", "Both")
    want_pptx = journal_fmt.value in ("PowerPoint (.pptx)", "Both")
    try:
        os.makedirs(folder, exist_ok=True)
        written = []
        if want_pptx:
            written.append(os.path.basename(
                _build_pptx(os.path.join(folder, base + ".pptx"))))
        if want_docx:
            written.append(os.path.basename(
                _build_docx(os.path.join(folder, base + ".docx"))))
    except ImportError as exc:
        export_status.object = (f"❌ Missing library: `{exc}`. Install with "
                                "`pip install --break-system-packages "
                                "python-docx python-pptx`.")
        return
    except Exception as exc:
        export_status.object = f"❌ Export failed: `{exc}`"
        return
    export_status.object = (f"✅ Exported {', '.join(written)} "
                            f"({len(JOURNAL)} items) to `{folder}`")


def _remove_last(event):
    if JOURNAL:
        JOURNAL.pop()
        _update_journal_status()
        export_status.object = f"Removed last — {len(JOURNAL)} remain."
    else:
        export_status.object = "Journal already empty."


def _clear_journal(event):
    JOURNAL.clear()
    _update_journal_status()
    export_status.object = "Journal cleared."


export_journal_btn.on_click(_export_journal)
remove_last_btn.on_click(_remove_last)
clear_journal_btn.on_click(_clear_journal)


def _refresh_listing(*events):
    base = os.path.expanduser(dir_input.value or "~")
    try:
        entries = os.listdir(base)
    except Exception as exc:
        subdir_select.options = ["—"]
        subdir_select.value = "—"
        file_select.options = []
        load_status.object = f"⚠️ Can't read `{base}`: `{exc}`"
        return
    dirs = sorted(f for f in entries
                  if os.path.isdir(os.path.join(base, f)) and not f.startswith("."))
    ims = sorted(f for f in entries if f.lower().endswith(".im"))
    subdir_select.options = ["—"] + dirs
    if subdir_select.value not in subdir_select.options:
        subdir_select.value = "—"
    file_select.options = ims


def _enter_subdir(event):
    if event.new and event.new != "—":
        dir_input.value = os.path.join(os.path.expanduser(dir_input.value), event.new)


def _go_up(event):
    cur = os.path.expanduser(dir_input.value or "~").rstrip("/")
    dir_input.value = os.path.dirname(cur) or "/"


def _load(event):
    if not file_select.value:
        load_status.object = "⚠️ Choose a `.im` file from the list."
        return
    path = os.path.join(os.path.expanduser(dir_input.value), file_select.value)
    load_status.object = f"Loading `{file_select.value}` …"
    try:
        img = MimsImage(path)
        state.img = img
        m = img.metadata
        load_status.object = (
            f"✅ **{file_select.value}**  \n"
            f"{m['width']}×{m['height']} px · {m['n_planes']} planes · "
            f"{m['n_masses']} channels · {m['field_um']:.2f} µm field"
        )
    except Exception as exc:
        state.img = None
        load_status.object = f"❌ Failed to load: `{exc}`"


dir_input.param.watch(_refresh_listing, "value")
subdir_select.param.watch(_enter_subdir, "value")
up_btn.on_click(_go_up)
load_btn.on_click(_load)
_refresh_listing()


# ─────────────────────────────────────────────────────────────────────────────
# Imaging tab
# ─────────────────────────────────────────────────────────────────────────────

def _render_channel(img, channel, cmap, upper_pct, sum_all, plane,
                    filt="None", strength=3.0, bin_n=1):
    """Single channel as a matplotlib figure. Binning (sum-pool, cropped) and
    filtering are applied to the DISPLAY array only — img.data is never touched."""
    stack = img.get_channel(channel)
    if sum_all:
        data = stack.sum(axis=0).astype(float)
        plane_lbl = "summed"
    else:
        p = int(min(max(plane, 0), stack.shape[0] - 1))
        data = stack[p].astype(float)
        plane_lbl = f"plane {p}"

    extra = []
    if bin_n > 1:
        data = _bin_sum(data, bin_n)
        extra.append(f"{bin_n}×{bin_n} bin")

    is_edge = filt in ("Sobel", "Canny")
    if filt == "Median":
        from scipy.ndimage import median_filter
        data = median_filter(data, size=max(1, int(strength)))
        extra.append(f"median {int(strength)}")
    elif filt == "Gaussian":
        from scipy.ndimage import gaussian_filter
        data = gaussian_filter(data, sigma=float(strength))
        extra.append(f"gaussian σ{strength:g}")
    elif filt == "Sobel":
        from skimage.filters import sobel
        data = sobel(data / (data.max() or 1.0))
        extra.append("sobel")
    elif filt == "Canny":
        from skimage.feature import canny
        data = canny(data / (data.max() or 1.0), sigma=float(strength)).astype(float)
        extra.append(f"canny σ{strength:g}")

    field = img.metadata["field_um"]
    if is_edge:
        vmin, vmax, cbar_label = 0.0, float(data.max() or 1.0), "edge"
    elif data.max() > 0:
        vmin, vmax, cbar_label = 0.0, float(np.percentile(data, upper_pct)), "counts"
    else:
        vmin, vmax, cbar_label = 0.0, 1.0, "counts"
    vmax = max(vmax, 1e-9)

    fig, ax = plt.subplots(figsize=(5.4, 5.4))
    im = ax.imshow(data, cmap=cmap, vmin=vmin, vmax=vmax,
                   extent=[0, field, field, 0], interpolation="nearest")
    source = "drift-corrected" if img.corrected is not None else "raw"
    suffix = (" — " + ", ".join(extra)) if extra else ""
    ax.set_title(f"{channel} — {plane_lbl}, {source}{suffix}", fontsize=10)
    ax.set_xlabel("µm")
    ax.set_ylabel("µm")
    fig.colorbar(im, ax=ax, fraction=0.046, pad=0.04, label=cbar_label)
    fig.tight_layout()
    plt.close(fig)
    return fig


def _build_fig(img, mode, ch, cm, pct, sa, pl, filt, strength, bin_n,
               nu, de, mc, mr, dr, hn, hd, inten, hcm, sc, rn, rx, hmed):
    if mode == "Single channel":
        return _render_channel(img, ch, cm, pct, sa, pl, filt, strength, bin_n)
    if mode == "Ratio":
        fig, _ = img.plot_ratio(nu, de, min_counts=(mc or None),
                                max_rel_err=(mr or None), delta_ref=(dr or None),
                                bin=bin_n, show=True)
        return fig
    fig, _ = img.plot_hsi(hn, hd, intensity=inten, cmap=hcm, scale_factor=sc,
                          ratio_min=(rn or None), ratio_max=(rx or None),
                          bin=bin_n, median_smooth=hmed, show=True)
    return fig


def imaging_view(img):
    if img is None:
        return pn.pane.Markdown("### Imaging\nLoad a `.im` file from the sidebar to begin.")

    masses = img.masses
    n_planes = int(img.metadata["n_planes"])
    non_se = [m for m in masses if "SE" not in m.upper()] or masses
    den_default = non_se[1] if len(non_se) > 1 else non_se[0]

    mode = pn.widgets.RadioButtonGroup(
        name="Mode", options=["Single channel", "Ratio", "HSI"],
        value="Single channel", button_type="primary")

    # Single-channel controls (bin is global, in the sidebar)
    channel = pn.widgets.Select(name="Channel", options=masses, value=masses[0], width=160)
    cmap = pn.widgets.Select(name="Colormap",
                             options=["gray", "viridis", "magma", "inferno", "cividis"],
                             value="gray", width=140)
    upper = pn.widgets.FloatSlider(name="Contrast (upper %ile)", start=90, end=100,
                                   step=0.1, value=99.5, width=230)
    sum_toggle = pn.widgets.Checkbox(name="Sum all planes", value=True)
    plane = pn.widgets.IntSlider(name="Plane", start=0,
                                 end=(n_planes - 1 if n_planes > 1 else 1),
                                 value=0, disabled=True, width=240)
    filt = pn.widgets.Select(name="Filter (display only)",
                             options=["None", "Median", "Gaussian", "Sobel", "Canny"],
                             value="None", width=170)
    strength = pn.widgets.FloatSlider(name="Filter size / σ", start=1, end=9,
                                      step=0.5, value=3, width=190)
    filt_note = pn.pane.Markdown(
        "<span style='font-size:11px;opacity:0.75'>Filters are **display only** — "
        "they do not touch or change the underlying data. Ratios, HSI and the "
        "Analysis tab always use the raw counts.</span>")

    def _toggle_plane(event):
        cur_n = int(img.metadata["n_planes"])
        plane.disabled = event.new or cur_n <= 1
    sum_toggle.param.watch(_toggle_plane, "value")

    # Ratio controls
    num = pn.widgets.Select(name="Numerator", options=masses, value=non_se[0], width=160)
    den = pn.widgets.Select(name="Denominator", options=masses, value=den_default, width=160)
    min_counts = pn.widgets.IntInput(name="Min counts (denom; 0 = off)", value=0, start=0, width=190)
    max_rel = pn.widgets.FloatInput(name="Max σ/R (0 = off)", value=0.0, start=0.0, step=0.05, width=150)
    dref = pn.widgets.FloatInput(name="δ reference (0 = image median)", value=0.0, start=0.0, width=210)

    # HSI controls
    hnum = pn.widgets.Select(name="Numerator", options=masses, value=non_se[0], width=150)
    hden = pn.widgets.Select(name="Denominator", options=masses, value=den_default, width=150)
    intensity = pn.widgets.Select(name="Intensity from",
                                  options=["denominator", "numerator", "sum"],
                                  value="denominator", width=150)
    hcmap = pn.widgets.Select(name="Colormap",
                              options=["viridis", "magma", "inferno", "hsv", "Classic OpenMIMS LUT"],
                              value="viridis", width=180)
    scale = pn.widgets.FloatInput(name="Scale factor", value=10000.0, start=1.0, width=120)
    rmin = pn.widgets.FloatInput(name="Hue min (0 = auto)", value=0.0, width=140)
    rmax = pn.widgets.FloatInput(name="Hue max (0 = auto)", value=0.0, width=140)
    hmed = pn.widgets.IntInput(name="Ratio median (px; 0/1 = off)", value=0, start=0, width=190)
    hsi_note = pn.pane.Markdown(
        "<span style='font-size:11px;opacity:0.75'>Ratio-median smooths the "
        "ratio **after** it is formed (display only). For genuine noise reduction "
        "use the global bin, which sum-pools the counts and enlarges the analytical "
        "volume.</span>")

    # ── Plane QC (run BEFORE drift correction) ────────────────────────────────
    qc_channel = pn.widgets.Select(name="Scan channel",
                                   options=["all channels"] + list(masses),
                                   value="all channels", width=170)
    qc_thresh = pn.widgets.FloatInput(name="Suspect threshold (% from median)",
                                      value=30.0, start=1.0, step=5.0, width=240)
    qc_scan_btn = pn.widgets.Button(name="Scan planes", button_type="default",
                                    icon="activity", width=130)
    qc_add_journal_btn = pn.widgets.Button(name="➕ Add scan to journal",
                                           button_type="default", width=190)
    qc_preview_btn = pn.widgets.Button(name="Auto-drop (preview)",
                                       button_type="default", width=170)
    qc_apply_btn = pn.widgets.Button(name="Apply auto-drop",
                                     button_type="warning", width=150)
    qc_manual = pn.widgets.TextInput(name="Drop planes (indices, e.g. 3,7,12)",
                                     value="", width=260)
    qc_add_current = pn.widgets.Button(name="↳ add scrubber plane",
                                       button_type="default", width=170)
    qc_drop_btn = pn.widgets.Button(name="Drop listed", button_type="warning",
                                    width=120)
    qc_reset_btn = pn.widgets.Button(name="↺ Reset to raw (reload)",
                                     button_type="default", width=200)
    qc_status = pn.pane.Markdown("")
    qc_plot_box = pn.Column(pn.pane.Markdown(
        "<span style='opacity:0.6'>No scan yet — click **Scan planes**.</span>"))
    qc_note = pn.pane.Markdown(
        "<span style='font-size:11px;opacity:0.75'>Scans per-plane total counts "
        "and flags outliers (charging events, glitches, blank frames). "
        "<b>Run before drift correction</b> — bad planes are the main cause of "
        "silent drift-correction failures. Dropping is destructive in memory; "
        "<b>Reset to raw</b> reloads the file from disk (this also clears drift). "
        "To eyeball a suspect plane first, switch to Single channel, untick "
        "“Sum all planes”, scrub to it, then “↳ add scrubber plane”. "
        "(plane_movie remains available in Colab for animated inspection.)</span>")

    def _qc_channel_arg():
        v = qc_channel.value
        return None if v == "all channels" else v

    def _run_scan():
        cur_n = int(img.metadata["n_planes"])
        if cur_n < 2:
            qc_plot_box.objects = [pn.pane.Markdown(
                "<span style='opacity:0.6'>Need ≥2 planes to scan.</span>")]
            return None
        try:
            res = img.plot_plane_diagnostics(
                channel=_qc_channel_arg(),
                threshold_pct=float(qc_thresh.value),
                show=False)
        except Exception as exc:
            qc_plot_box.objects = [pn.pane.Alert(f"Scan failed: `{exc}`",
                                                 alert_type="danger")]
            return None
        fig = res["figure"]
        pane = pn.pane.Matplotlib(fig, dpi=100, tight=True, height=300)
        plt.close(fig)
        qc_plot_box.objects = [pane]
        return res

    def _after_inplace_change():
        """Refresh slider bounds, drift status, the main figure and the scan
        after an in-place plane drop (img.data was mutated under us)."""
        cur_n = int(img.metadata["n_planes"])
        new_end = cur_n - 1 if cur_n > 1 else 1
        plane.end = new_end
        if plane.value > new_end:
            plane.value = new_end
        plane.disabled = sum_toggle.value or cur_n <= 1
        drift_status.object = ("Drift: **applied**" if img.corrected is not None
                               else "Drift: not applied")
        state.plane_version += 1     # force the main image to re-render
        _run_scan()                  # keep the diagnostics plot current

    def _on_scan(event):
        res = _run_scan()
        if res is None:
            return
        sus = res["suspect_indices"]
        t = float(qc_thresh.value)
        if sus:
            qc_status.object = (f"Scan: **{len(sus)}** suspect plane(s) at "
                                f"±{t:.0f}%: `{sus}` · "
                                f"{int(img.metadata['n_planes'])} planes total.")
        else:
            qc_status.object = (f"Scan: no suspects at ±{t:.0f}% · "
                                f"{int(img.metadata['n_planes'])} planes.")
    qc_scan_btn.on_click(_on_scan)

    def _add_scan_to_journal(event):
        cur_n = int(img.metadata["n_planes"])
        if cur_n < 2:
            qc_status.object = "Need ≥2 planes to scan."
            return
        try:
            res = img.plot_plane_diagnostics(
                channel=_qc_channel_arg(),
                threshold_pct=float(qc_thresh.value),
                show=False)
        except Exception as exc:
            qc_status.object = f"❌ Couldn't render scan: `{exc}`"
            return
        fig = res["figure"]
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                    facecolor=fig.get_facecolor())
        plt.close(fig)
        sus = res["suspect_indices"]
        chan = qc_channel.value
        cap = (f"Plane QC scan — {chan}, ±{float(qc_thresh.value):.0f}% threshold, "
               f"{len(sus)} of {cur_n} flagged  ·  {os.path.basename(img.path)}")
        JOURNAL.append({
            "png": buf.getvalue(),
            "caption": cap,
            "added": datetime.datetime.now().strftime("%Y-%m-%d %H:%M"),
        })
        _update_journal_status()
        qc_status.object = f"✅ Scan added to journal — {len(JOURNAL)} view(s)."
    qc_add_journal_btn.on_click(_add_scan_to_journal)

    def _on_preview(event):
        cur_n = int(img.metadata["n_planes"])
        t = float(qc_thresh.value)
        if cur_n < 2:
            qc_status.object = "Need ≥2 planes to scan."
            return
        try:
            bad = img.auto_drop_bad_planes(threshold_pct=t,
                                           channel=_qc_channel_arg(),
                                           dry_run=True, verbose=False)
        except Exception as exc:
            qc_status.object = f"Preview failed: `{exc}`"
            return
        if bad:
            qc_status.object = (f"Auto-drop **preview**: would drop {len(bad)} "
                                f"plane(s) `{bad}` → {cur_n - len(bad)} remain. "
                                f"Click **Apply auto-drop** to commit.")
        else:
            qc_status.object = f"Auto-drop preview: nothing flagged at ±{t:.0f}%."
    qc_preview_btn.on_click(_on_preview)

    def _on_apply(event):
        cur_n = int(img.metadata["n_planes"])
        t = float(qc_thresh.value)
        if cur_n < 2:
            qc_status.object = "Need ≥2 planes to scan."
            return
        try:
            dropped = img.auto_drop_bad_planes(threshold_pct=t,
                                               channel=_qc_channel_arg(),
                                               dry_run=False, verbose=False)
        except Exception as exc:
            qc_status.object = f"Auto-drop failed: `{exc}`"
            return
        if dropped:
            _after_inplace_change()
            qc_status.object = (f"✅ Dropped {len(dropped)} plane(s) `{dropped}` "
                                f"→ {int(img.metadata['n_planes'])} planes remain. "
                                f"(Reset to raw to undo.)")
        else:
            qc_status.object = f"Nothing flagged at ±{t:.0f}% — no change."
    qc_apply_btn.on_click(_on_apply)

    def _parse_indices(text):
        out = []
        for tok in text.replace(";", ",").split(","):
            tok = tok.strip()
            if tok:
                out.append(int(tok))      # ValueError propagates to caller
        return out

    def _on_add_current(event):
        cur = qc_manual.value.strip()
        add = str(int(plane.value))
        qc_manual.value = (cur + ", " + add) if cur else add
    qc_add_current.on_click(_on_add_current)

    def _on_drop_listed(event):
        try:
            idx = _parse_indices(qc_manual.value)
        except ValueError:
            qc_status.object = ("⚠️ Couldn't parse indices — use comma-separated "
                                "integers, e.g. `3,7,12`.")
            return
        if not idx:
            qc_status.object = "⚠️ No plane indices entered."
            return
        try:
            dropped = img.drop_planes(idx, verbose=False)
        except Exception as exc:
            qc_status.object = f"❌ Drop failed: `{exc}`"
            return
        if dropped:
            _after_inplace_change()
            qc_manual.value = ""
            qc_status.object = (f"✅ Dropped {len(dropped)} plane(s) `{dropped}` "
                                f"→ {int(img.metadata['n_planes'])} planes remain. "
                                f"(Reset to raw to undo.)")
        else:
            qc_status.object = "No valid planes dropped (out of range?)."
    qc_drop_btn.on_click(_on_drop_listed)

    def _on_reset(event):
        try:
            fresh = MimsImage(img.path)
        except Exception as exc:
            qc_status.object = f"❌ Reset failed: `{exc}`"
            return
        # New object → the tab rebuilds from disk-fresh data (drift + drops cleared).
        state.img = fresh
    qc_reset_btn.on_click(_on_reset)

    qc_section = pn.Column(
        pn.pane.Markdown("**Plane QC**  <span style='font-size:11px;opacity:0.7'>"
                         "(run before drift correction)</span>"),
        pn.Row(qc_channel, qc_thresh, qc_scan_btn, qc_add_journal_btn),
        pn.Row(qc_preview_btn, qc_apply_btn),
        pn.Row(qc_manual, qc_add_current, qc_drop_btn),
        pn.Row(qc_reset_btn),
        qc_status,
        qc_plot_box,
        qc_note,
    )

    # Drift correction (common)
    se_default = "SE" if any("SE" in m.upper() for m in masses) else masses[0]
    drift_ref = pn.widgets.Select(name="Drift reference", options=masses,
                                  value=se_default, width=170)
    drift_btn = pn.widgets.Button(name="Run drift correction",
                                  button_type="default", icon="arrows-move")
    drift_status = pn.pane.Markdown(
        "Drift: **applied**" if img.corrected is not None else "Drift: not applied")

    def _do_drift(event):
        try:
            img.drift_correct(reference=drift_ref.value)
            drift_status.object = "Drift: **applied**"
            state.drift_version += 1
        except Exception as exc:
            drift_status.object = f"Drift failed: `{exc}`"
    drift_btn.on_click(_do_drift)

    # Save current view at chosen DPI
    out_dir = pn.widgets.TextInput(name="Save to (folder)",
                                   value=os.path.expanduser("~"), width=240)
    save_name = pn.widgets.TextInput(name="Filename", value="figure", width=160)
    fmt = pn.widgets.Select(name="Format", options=["png", "pdf", "svg", "tiff"],
                            value="png", width=90)
    dpi_in = pn.widgets.IntInput(name="DPI", value=600, start=72, end=1200, width=90)
    save_btn = pn.widgets.Button(name="💾 Save view", button_type="default", width=120)
    save_status = pn.pane.Markdown("")

    def _save(event):
        try:
            fig = _build_fig(img, mode.value, channel.value, cmap.value, upper.value,
                             sum_toggle.value, plane.value, filt.value, strength.value,
                             state.bin_n, num.value, den.value, min_counts.value,
                             max_rel.value, dref.value, hnum.value, hden.value,
                             intensity.value, hcmap.value, scale.value, rmin.value,
                             rmax.value, hmed.value)
            folder = os.path.expanduser(out_dir.value or "~")
            os.makedirs(folder, exist_ok=True)
            path = os.path.join(folder, f"{save_name.value or 'figure'}.{fmt.value}")
            save_figure(fig, path, dpi=int(dpi_in.value))
            plt.close(fig)
            save_status.object = f"✅ Saved `{path}` at {int(dpi_in.value)} DPI"
        except Exception as exc:
            save_status.object = f"❌ Save failed: `{exc}`"
    save_btn.on_click(_save)

    # Add the current view to the journal (captured as a PNG with an auto caption)
    add_journal_btn = pn.widgets.Button(name="➕ Add current view to journal",
                                        button_type="default", width=240)
    journal_local = pn.pane.Markdown("")

    def _current_caption():
        src = "drift-corrected" if img.corrected is not None else "raw"
        fname = os.path.basename(img.path)
        binlbl = (f", {state.bin_n}×{state.bin_n} bin" if state.bin_n > 1 else "")
        m = mode.value
        if m == "Single channel":
            plane_lbl = "summed" if sum_toggle.value else f"plane {plane.value}"
            flt = ("" if filt.value == "None"
                   else f", {filt.value.lower()} {strength.value:g}")
            return (f"{channel.value} — {plane_lbl}, {src}{binlbl}{flt}"
                    f"  ·  {fname}")
        if m == "Ratio":
            return f"{num.value} / {den.value} ratio, {src}{binlbl}  ·  {fname}"
        return (f"{hnum.value} / {hden.value} HSI "
                f"(intensity={intensity.value}), {src}{binlbl}  ·  {fname}")

    def _render_current_png():
        fig = _build_fig(img, mode.value, channel.value, cmap.value, upper.value,
                         sum_toggle.value, plane.value, filt.value, strength.value,
                         state.bin_n, num.value, den.value, min_counts.value,
                         max_rel.value, dref.value, hnum.value, hden.value,
                         intensity.value, hcmap.value, scale.value, rmin.value,
                         rmax.value, hmed.value)
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
        plt.close(fig)
        return buf.getvalue()

    def _add_to_journal(event):
        try:
            png = _render_current_png()
        except Exception as exc:
            journal_local.object = f"❌ Couldn't render view: `{exc}`"
            return
        JOURNAL.append({
            "png": png,
            "caption": _current_caption(),
            "added": datetime.datetime.now().strftime("%Y-%m-%d %H:%M"),
        })
        _update_journal_status()
        journal_local.object = f"✅ Added — {len(JOURNAL)} view(s) in journal."
    add_journal_btn.on_click(_add_to_journal)

    def _controls(m):
        if m == "Single channel":
            return pn.Column(pn.Row(channel, cmap, upper),
                             pn.Row(sum_toggle, plane),
                             pn.Row(filt, strength),
                             filt_note)
        if m == "Ratio":
            return pn.Column(pn.Row(num, den), pn.Row(min_counts, max_rel, dref))
        return pn.Column(pn.Row(hnum, hden, intensity),
                         pn.Row(hcmap, scale, rmin, rmax),
                         pn.Row(hmed), hsi_note)

    def _view(m, ch, cm, pct, sa, pl, fl, st, nu, de, mc, mr, dr,
              hn, hd, inten, hcm, sc, rn, rx, hm, bn, _ver, _pver):
        try:
            fig = _build_fig(img, m, ch, cm, pct, sa, pl, fl, st, bn,
                             nu, de, mc, mr, dr, hn, hd, inten, hcm, sc, rn, rx, hm)
        except Exception as exc:
            return pn.pane.Alert(f"Render failed: `{exc}`", alert_type="danger")
        pane = pn.pane.Matplotlib(fig, dpi=110, tight=True, height=480)
        plt.close(fig)
        return pane

    view = pn.bind(_view, mode.param.value, channel.param.value, cmap.param.value,
                   upper.param.value, sum_toggle.param.value, plane.param.value,
                   filt.param.value, strength.param.value,
                   num.param.value, den.param.value, min_counts.param.value,
                   max_rel.param.value, dref.param.value, hnum.param.value,
                   hden.param.value, intensity.param.value, hcmap.param.value,
                   scale.param.value, rmin.param.value, rmax.param.value,
                   hmed.param.value, state.param.bin_n, state.param.drift_version,
                   state.param.plane_version)

    return pn.Column(
        mode,
        pn.panel(pn.bind(_controls, mode.param.value)),
        pn.layout.Divider(),
        qc_section,
        pn.layout.Divider(),
        pn.Row(drift_ref, drift_btn, drift_status),
        pn.Row(out_dir, save_name, fmt, dpi_in, save_btn),
        save_status,
        pn.Row(add_journal_btn, journal_local),
        pn.panel(view),
    )


# ─────────────────────────────────────────────────────────────────────────────
# Metadata tab  (working)
# ─────────────────────────────────────────────────────────────────────────────

def metadata_view(img):
    if img is None:
        return pn.pane.Markdown("### Metadata\nLoad a `.im` file to view the full Cameca header.")
    try:
        meta = img.read_full_metadata()
    except Exception as exc:
        return pn.pane.Alert(f"Could not read metadata: `{exc}`", alert_type="danger")

    panes = []
    layout = meta.get("layout", {})
    if not layout.get("layout_supported", True):
        reasons = "\n".join(f"- {r}" for r in layout.get("reasons", []))
        panes.append(pn.pane.Alert(
            "**Header layout not recognised** — the analytical values below may be "
            "unreliable for this file.\n\n" + reasons, alert_type="warning"))
    else:
        panes.append(pn.pane.Alert(
            "Header layout recognised — values validated for this acquisition type.",
            alert_type="success"))

    text = pymims_metadata.format_metadata(meta, masses=img.masses)
    panes.append(pn.pane.Str(text, styles={
        "font-family": "ui-monospace, SFMono-Regular, Menlo, monospace",
        "white-space": "pre-wrap",
        "font-size": "12px",
        "line-height": "1.35",
    }))
    return pn.Column(*panes, sizing_mode="stretch_width")


# ─────────────────────────────────────────────────────────────────────────────
# HMR + Analysis tabs  (placeholders)
# ─────────────────────────────────────────────────────────────────────────────

def hmr_view():
    if not _HMR_OK:
        return pn.pane.Markdown(
            "### HMR — high mass resolution\n\n"
            "The HMR engine modules aren't importable. Drop `irf.py`, "
            "`deconvolve.py`, `sparse_deconvolve.py` and `hmr_identify.py` in "
            "the repo root (next to this app), then restart `panel serve`.\n\n"
            f"<span style='opacity:0.6'>Import error: `{_HMR_ERR}`</span>")

    import os as _os
    H = {"m": None, "counts": None, "axes": {}, "axis_label": {}, "path": None,
         "fit": None, "fit_axis": None, "ident_labels": None, "nominal": None,
         "ident_text": None, "ident_caption": None}

    # ── helpers ──────────────────────────────────────────────────────────────
    def _load_hmr(path):
        import pandas as pd
        low = path.lower()
        if low.endswith((".xlsx", ".xls")):
            try:
                df = pd.read_excel(path)
            except Exception:
                # openpyxl (or the .xls engine) missing/failed — for .xlsx we can
                # still read it with the stdlib fallback; .xls has no fallback.
                if low.endswith(".xlsx"):
                    df = _read_xlsx_minimal(path)
                else:
                    raise
        else:
            df = pd.read_csv(path)
        cols = list(df.columns)

        def find(names):
            for c in cols:
                cl = (str(c).lower().replace("(", "").replace(")", "")
                      .replace("µ", "u").replace(" ", "").strip())
                if cl in names:
                    return c
            return None
        x_v = find({"xv"})
        x_amu = find({"xamu", "mass", "amu"})
        x_um = find({"xum"})
        axis_cols = {c for c in (x_v, x_amu, x_um) if c}
        # counts = non-axis, non-Pt numeric column with the largest range
        best, best_rng = None, -1.0
        for c in cols:
            if c in axis_cols or str(c).lower() == "pt":
                continue
            v = pd.to_numeric(df[c], errors="coerce")
            if v.notna().sum() < 3:
                continue
            v = v.fillna(0.0)
            rng = float(v.max() - v.min())
            if rng > best_rng:
                best, best_rng = c, rng
        if best is None:
            raise ValueError("No counts column found.")
        counts = pd.to_numeric(df[best], errors="coerce").fillna(0.0).to_numpy(float).copy()
        counts[counts < 1e-3] = 0.0                      # 1e-6 baseline floor -> 0
        axes, labels = {}, {}
        if x_amu is not None:
            axes["amu"] = df[x_amu].to_numpy(float); labels["amu"] = "mass (amu, nominal)"
        if x_v is not None:
            axes["V"] = df[x_v].to_numpy(float); labels["V"] = "deflector (V)"
        if x_um is not None:
            axes["µm"] = df[x_um].to_numpy(float); labels["µm"] = "position (µm)"
        if not axes:                                     # fall back to point index
            axes["pt"] = np.arange(len(counts), dtype=float); labels["pt"] = "point"
        return axes, labels, counts, best, len(counts)

    def _mrp_10_90(m, counts, side="rising", m_ref=None):
        """Resolving power from the 10–90% edge width. ΔM = |m90 − m10|;
        MRP = M / ΔM with M defaulting to the peak-centre value on this axis."""
        m = np.asarray(m, float); c = np.asarray(counts, float)
        pk = int(np.argmax(c))
        base = float(np.median(c[: max(3, len(c) // 10)]))
        plateau = float(np.median(c[max(0, pk - 2):pk + 3]))
        amp = max(plateau - base, 1e-9)
        t10, t90 = base + 0.10 * amp, base + 0.90 * amp
        seg = slice(0, pk + 1) if side == "rising" else slice(pk, len(c))
        xs, ys = m[seg], c[seg]

        def cross(level):
            for i in range(1, len(ys)):
                if (ys[i - 1] - level) * (ys[i] - level) <= 0 and ys[i] != ys[i - 1]:
                    t = (level - ys[i - 1]) / (ys[i] - ys[i - 1])
                    return xs[i - 1] + t * (xs[i] - xs[i - 1])
            return None
        m10, m90 = cross(t10), cross(t90)
        if m10 is None or m90 is None:
            return None
        dM = abs(m90 - m10)
        center = float(m[pk])
        M = float(m_ref) if m_ref else abs(round(center)) or abs(center)
        mrp = (M / dM) if dM > 0 else float("nan")
        return dict(dM=dM, mrp=mrp, m10=m10, m90=m90, center=center, M=M,
                    plateau=plateau, base=base)

    # ── widgets: load ────────────────────────────────────────────────────────
    hmr_dir = pn.widgets.TextInput(name="HMR folder",
                                   value=_os.path.expanduser("~"), width=320)
    hmr_file = pn.widgets.Select(name="HMR file", options=[], width=320)
    hmr_scan_btn = pn.widgets.Button(name="↻ Scan folder", button_type="default",
                                     width=130)
    hmr_load_btn = pn.widgets.Button(name="⤓ Load scan", button_type="primary",
                                     width=130)
    hmr_info = pn.pane.Markdown("<span style='opacity:0.6'>No scan loaded.</span>")

    def _scan_folder(event=None):
        try:
            files = sorted(f for f in _os.listdir(_os.path.expanduser(hmr_dir.value))
                           if f.lower().endswith((".csv", ".xlsx", ".xls", ".txt", ".hmr")))
        except Exception as exc:
            hmr_info.object = f"❌ Can't list folder: `{exc}`"
            return
        hmr_file.options = files
        if files:
            hmr_file.value = files[0]
    hmr_scan_btn.on_click(_scan_folder)

    # ── widgets: spectrum ──────────────────────────────────────────────────────
    x_axis = pn.widgets.Select(name="X axis", options=["amu"], value="amu", width=130)
    log_y = pn.widgets.Checkbox(name="Log Y", value=True)
    spectrum_box = pn.Column(pn.pane.Markdown(
        "<span style='opacity:0.6'>Load a scan to see its spectrum.</span>"))

    # ── widgets: IRF / MRP ─────────────────────────────────────────────────────
    edge_side = pn.widgets.Select(name="Edge", options=["rising", "falling"],
                                  value="rising", width=120)
    m_ref = pn.widgets.FloatInput(name="M (override, optional)", value=0.0,
                                  width=180)
    mrp_btn = pn.widgets.Button(name="Compute IRF + MRP", button_type="primary",
                                width=190)
    mrp_out = pn.pane.Markdown("")

    # ── widgets: deconvolution ─────────────────────────────────────────────────
    dc_method = pn.widgets.RadioButtonGroup(
        name="Method", options=["Sparse (auto-N)", "N-sweep"],
        value="Sparse (auto-N)", button_type="primary")
    dc_nmin = pn.widgets.IntInput(name="N min", value=1, start=1, end=8, width=90)
    dc_nmax = pn.widgets.IntInput(name="N max", value=4, start=1, end=10, width=90)
    dc_grid = pn.widgets.FloatInput(name="Grid step (mau)", value=0.5, start=0.1,
                                    step=0.1, width=130)
    dc_stab = pn.widgets.FloatSlider(name="Stability threshold", start=0.02,
                                     end=0.5, step=0.02, value=0.10, width=200)
    dc_minsnr = pn.widgets.FloatInput(name="Min peak a/σ", value=3.0, start=0.0,
                                      end=20.0, step=0.5, width=120)
    dc_run = pn.widgets.Button(name="Run deconvolution", button_type="primary",
                               width=190)
    dc_status = pn.pane.Markdown("")
    dc_table = pn.Column()
    resid_box = pn.Column()

    # ── widgets: species ID ────────────────────────────────────────────────────
    _sel_elems = set(["C", "N", "H", "O"])           # default selection
    _elem_pos = {(e["row"], e["col"]): e["sym"] for e in HMR_ELEMENTS}
    _elem_toggles = {}
    sel_label = pn.pane.Markdown("")

    def _refresh_sel_label():
        sel_label.object = ("**Selected:** "
                            + (", ".join(sorted(_sel_elems)) or "_none_"))

    def _make_toggle(sym):
        tog = pn.widgets.Toggle(name=sym, value=(sym in _sel_elems),
                                width=34, height=28, margin=(1, 1))

        def _cb(event, _sym=sym):
            (_sel_elems.add if event.new else _sel_elems.discard)(_sym)
            _refresh_sel_label()
            _refresh_ref_species()
        tog.param.watch(_cb, "value")
        _elem_toggles[sym] = tog
        return tog

    _pt_layout = [1, 2, 3, 4, 5, 6, None, 9, 10]  # None = lanthanide/actinide gap
    _pt_rows = []
    for _r in _pt_layout:
        if _r is None:
            _pt_rows.append(pn.Spacer(height=6))
            continue
        _cells = []
        for _c in range(1, 19):
            _sym = _elem_pos.get((_r, _c))
            _cells.append(_make_toggle(_sym) if _sym
                          else pn.Spacer(width=34, height=28, margin=(1, 1)))
        _pt_rows.append(pn.Row(*_cells, margin=0))
    periodic_table = pn.Column(*_pt_rows, margin=(4, 0))

    id_clear = pn.widgets.Button(name="Clear", button_type="default", width=90)

    def _clear_elems(event=None):
        _sel_elems.clear()
        for t in _elem_toggles.values():
            t.value = False
        _refresh_sel_label()
        _refresh_ref_species()
    id_clear.on_click(_clear_elems)
    _refresh_sel_label()

    id_pol = pn.widgets.Select(name="Polarity", options=["anion (−)", "cation (+)"],
                               value="anion (−)", width=140)
    id_maxatoms = pn.widgets.IntInput(name="Max atoms", value=4, start=1, end=6,
                                      width=110)
    id_labeled = pn.widgets.Checkbox(
        name="Labeled sample — ignore natural abundances", value=False)
    id_mode = pn.widgets.RadioButtonGroup(
        name="Mode", options=["Reference peak", "Blind (best guess)"],
        value="Reference peak", button_type="default")
    id_ref_peak = pn.widgets.Select(name="Reference peak", options=[], width=240)
    id_ref_species = pn.widgets.Select(name="is species", options=[], width=180)
    id_btn = pn.widgets.Button(name="Identify peaks", button_type="primary",
                               width=160)
    id_out = pn.Column()
    id_ref_row = pn.Row(id_ref_peak, id_ref_species)

    def _current_candidates():
        elems = sorted(_sel_elems)
        if not elems or H["m"] is None:
            return []
        pol = -1 if id_pol.value.startswith("anion") else 1
        nominal = H.get("nominal") or 0
        return enumerate_candidates(
            elems, nominal, max_atoms=int(id_maxatoms.value), polarity=pol,
            mz_tol=0.1, ignore_abundance=bool(id_labeled.value))

    def _refresh_ref_species(event=None):
        cands = _current_candidates()
        opts = [c.formula for c in cands]            # already plausibility-sorted
        id_ref_species.options = opts
        if opts and id_ref_species.value not in opts:
            id_ref_species.value = opts[0]

    def _toggle_ref_row(event=None):
        id_ref_row.visible = (id_mode.value == "Reference peak")
    id_mode.param.watch(_toggle_ref_row, "value")
    for _w in (id_pol, id_maxatoms, id_labeled):
        _w.param.watch(_refresh_ref_species, "value")

    # ── plotting ───────────────────────────────────────────────────────────────
    def _build_spectrum_fig():
        ax_key = x_axis.value
        if H["m"] is None or ax_key not in H["axes"]:
            return None
        m = H["axes"][ax_key]
        y = H["counts"]
        fig, ax = plt.subplots(figsize=(9.2, 4.7))
        ax.plot(m, y, lw=1.0, color="#1f77b4", label="data", zorder=3)
        # model overlay (only meaningful on the axis the fit was run on)
        res = H["fit"]
        comp_colors = []
        if res is not None and res.n_peaks > 0 and H["fit_axis"] == ax_key:
            mm = np.linspace(float(m.min()), float(m.max()), 1200)
            ax.plot(mm, res.model(mm), "k--", lw=1.3, label="model", zorder=4)
            cyc = plt.rcParams["axes.prop_cycle"].by_key().get(
                "color", ["#ff7f0e", "#2ca02c", "#9467bd", "#8c564b"])
            for j, p in enumerate(res.peaks):
                col = cyc[j % len(cyc)]
                comp_colors.append(col)
                comp = rect_gauss(mm, p.a, p.mu, res.w, res.sigma)
                ax.plot(mm, comp, lw=0.9, alpha=0.7, color=col, zorder=2)
                ax.axvline(p.mu, color=col, ls=":", lw=0.9, zorder=1)
        # species-ID labels, coloured to match each peak's component
        labels = H.get("ident_labels")
        if labels and H["fit_axis"] == ax_key:
            ymax = max(float(np.max(y)), 1.0)
            for j, (mo, formula) in enumerate(labels):
                col = comp_colors[j] if j < len(comp_colors) else "#2ca02c"
                ax.annotate(formula, xy=(mo, ymax), xytext=(0, 6),
                            textcoords="offset points", fontsize=8, rotation=90,
                            ha="center", va="bottom", color=col)
        if log_y.value:
            ax.set_yscale("log")
            pos = y[y > 0]
            ax.set_ylim(bottom=max(0.5, float(pos.min()) * 0.5) if pos.size else 0.5)
        ax.set_xlabel(H["axis_label"].get(ax_key, ax_key))
        ax.set_ylabel("counts")
        ax.set_title(f"HMR — {_os.path.basename(H['path'])}", fontsize=10)
        ax.legend(fontsize=8, loc="upper right")
        fig.tight_layout()
        return fig

    def _render_spectrum():
        fig = _build_spectrum_fig()
        if fig is None:
            return
        spectrum_box.objects = [pn.pane.Matplotlib(fig, dpi=110, tight=True,
                                                   sizing_mode="stretch_width")]

    def _build_residuals_fig():
        res = H["fit"]
        ax_key = H["fit_axis"]
        if res is None or ax_key is None or ax_key not in H["axes"]:
            return None
        m = H["axes"][ax_key]
        fig, ax = plt.subplots(figsize=(9.2, 2.4))
        ax.axhline(0, color="#888", lw=0.8)
        ax.plot(m, res.residuals, lw=1.0, color="#9467bd")
        ax.set_xlabel(H["axis_label"].get(ax_key, ax_key))
        ax.set_ylabel("residual\n(data − model)")
        ax.set_title("Deconvolution residuals", fontsize=9)
        fig.tight_layout()
        return fig

    def _load(event=None):
        if not hmr_file.value:
            hmr_info.object = "⚠️ Pick a file (Scan folder first)."
            return
        path = _os.path.join(_os.path.expanduser(hmr_dir.value), hmr_file.value)
        try:
            axes, labels, counts, ccol, n = _load_hmr(path)
        except Exception as exc:
            hmr_info.object = f"❌ Load failed: `{exc}`"
            return
        _m_axis = axes.get("amu", next(iter(axes.values())))
        H.update(path=path, axes=axes, axis_label=labels, counts=counts,
                 m=_m_axis, fit=None, fit_axis=None, ident_labels=None,
                 ident_text=None, ident_caption=None,
                 nominal=int(round(float(np.median(_m_axis)))))
        HMR_SHARED.update(m=_m_axis, counts=counts, axes=axes,
                          axis_label=labels, name=_os.path.basename(path))
        id_ref_peak.options = []; id_ref_species.options = []
        x_axis.options = list(axes.keys())
        x_axis.value = "amu" if "amu" in axes else list(axes.keys())[0]
        amu = axes.get("amu")
        rng = (f"{amu.min():.4f}–{amu.max():.4f} amu" if amu is not None
               else f"{n} pts")
        hmr_info.object = (f"✅ **{_os.path.basename(path)}** · {n} pts · "
                           f"{rng} · counts in `{ccol}` (max {counts.max():.0f})")
        dc_status.object = ""; dc_table.objects = []; resid_box.objects = []
        id_out.objects = []; mrp_out.object = ""
        _render_spectrum()
    hmr_load_btn.on_click(_load)
    x_axis.param.watch(lambda e: _render_spectrum(), "value")
    log_y.param.watch(lambda e: _render_spectrum(), "value")

    def _compute_mrp(event=None):
        if H["m"] is None:
            mrp_out.object = "⚠️ Load a scan first."
            return
        ax_key = x_axis.value
        m = H["axes"][ax_key]
        r = _mrp_10_90(m, H["counts"], side=edge_side.value,
                       m_ref=(m_ref.value or None))
        if r is None:
            mrp_out.object = "❌ Couldn't find 10%/90% crossings on this edge."
            return
        try:
            sigma, info = fit_sigma_from_edge(m, H["counts"], side=edge_side.value)
        except Exception:
            sigma = float("nan")
        unit = ax_key
        mrp_out.object = (
            f"**MRP (10–90% {edge_side.value} edge)** = "
            f"**{r['mrp']:.0f}**  ·  ΔM = {r['dM']*1000:.3f} m{unit} "
            f"(M = {r['M']:.4g} {unit})  ·  edge σ = {sigma*1000:.3f} m{unit}  ·  "
            f"plateau {r['plateau']:.0f}, baseline {r['base']:.0f}")
    mrp_btn.on_click(_compute_mrp)

    def _run_deconv(event=None):
        if H["m"] is None:
            dc_status.object = "⚠️ Load a scan first."
            return
        ax_key = x_axis.value
        m = H["axes"][ax_key]; y = H["counts"]
        if float(np.max(y)) <= 0:
            dc_status.object = "⚠️ No counts above baseline in this scan."
            return
        dc_status.object = "Running…"
        try:
            if dc_method.value.startswith("Sparse"):
                res = deconvolve_sparse(
                    m, y, grid_step_mau=float(dc_grid.value),
                    stability_threshold=float(dc_stab.value), verbose=False)
                method_lbl = "sparse / stability-selection"
            else:
                sweep = deconvolve_with_N_sweep(
                    m, y, N_range=(int(dc_nmin.value), int(dc_nmax.value)),
                    verbose=False)
                if not sweep:
                    dc_status.object = "❌ N-sweep produced no fits."
                    return
                res = min(sweep, key=lambda r: r.bic)
                method_lbl = f"N-sweep, BIC-selected (N={res.n_peaks})"
        except Exception as exc:
            dc_status.object = f"❌ Deconvolution failed: `{exc}`"
            return
        # cull statistically insignificant peaks (a/σ below threshold) and
        # refit the survivors with the peak shape fixed — removes the null,
        # near-degenerate satellites the fitter parks on resolved peaks.
        n_culled = 0
        sig = float(dc_minsnr.value)
        if sig > 0 and res.n_peaks > 1:
            keep = [p for p in res.peaks
                    if p.a_se > 0 and (p.a / p.a_se) >= sig]
            n_culled = res.n_peaks - len(keep)
            if keep and n_culled > 0:
                try:
                    res2 = deconvolve(
                        m, y, N=len(keep),
                        initial_positions=np.array([p.mu for p in keep]),
                        sigma0=res.sigma, w0=res.w, fix_sigma=True, fix_w=True,
                        verbose=False)
                    res = res2
                except Exception:
                    n_culled = 0          # refit failed: keep the original fit
        H["fit"] = res; H["fit_axis"] = ax_key
        H["ident_labels"] = None; H["ident_text"] = None; H["ident_caption"] = None
        # populate species-ID reference-peak choices from the fitted peaks
        _peak_opts = {f"Peak {i+1} @ {p.mu:.5f}": i
                      for i, p in enumerate(res.peaks)}
        id_ref_peak.options = _peak_opts
        if _peak_opts:
            id_ref_peak.value = list(_peak_opts.values())[0]
        _refresh_ref_species()
        # peak table
        rows = []
        anchor = res.peaks[0].mu if res.peaks else 0.0
        for i, p in enumerate(res.peaks):
            rows.append([f"{i+1}", f"{p.mu:.5f}", f"{p.a:.0f}",
                         f"±{p.a_se:.0f}", f"{(p.mu-anchor)*1000:+.2f}"])
        header = ["#", f"centre ({ax_key})", "amp", "amp ±", f"Δ from #1 (m{ax_key})"]
        dc_table.objects = [_html_table(header, rows, width=600)]
        dc_status.object = (
            f"✅ {method_lbl}: **{res.n_peaks} peak(s)** · "
            f"w={res.w*1000:.3f} m{ax_key}, σ={res.sigma*1000:.3f} m{ax_key} · "
            f"χ²={res.chi2:.1f}, BIC={res.bic:.1f}, AIC={res.aic:.1f}"
            + (f" · culled {n_culled} below a/σ={sig:g} and refit"
               if n_culled else ""))
        # residuals plot
        fig = _build_residuals_fig()
        if fig is not None:
            resid_box.objects = [pn.pane.Matplotlib(fig, dpi=110, tight=True,
                                                    sizing_mode="stretch_width")]
        _render_spectrum()
    dc_run.on_click(_run_deconv)

    def _identify(event=None):
        res = H["fit"]
        if res is None or res.n_peaks == 0:
            id_out.objects = [pn.pane.Markdown(
                "⚠️ Run a deconvolution first — identification works on the "
                "fitted peak centres.")]
            return
        elems = sorted(_sel_elems)
        if not elems:
            id_out.objects = [pn.pane.Markdown(
                "⚠️ Select at least one element on the periodic table.")]
            return
        pol = -1 if id_pol.value.startswith("anion") else 1
        peaks_mz = [p.mu for p in res.peaks]
        labeled = bool(id_labeled.value)
        maxa = int(id_maxatoms.value)
        H["ident_labels"] = None

        if id_mode.value == "Reference peak":
            ref_i = id_ref_peak.value
            ref_f = id_ref_species.value
            if ref_i is None or not ref_f:
                id_out.objects = [pn.pane.Markdown(
                    "⚠️ Pick a reference peak and the species it is.")]
                return
            try:
                res_id = identify_by_difference(
                    peaks_mz, int(ref_i), ref_f, elems, polarity=pol,
                    max_atoms=maxa, nominal_mass=H.get("nominal"),
                    ignore_abundance=labeled)
            except Exception as exc:
                id_out.objects = [pn.pane.Markdown(f"❌ Identify failed: `{exc}`")]
                return
            if res_id is None:
                id_out.objects = [pn.pane.Markdown(
                    f"⚠️ `{ref_f}` isn't among the candidates for the selected "
                    "elements — widen Max atoms or re-check the selection.")]
                return
            rows = []
            for a in res_id:
                mark = " ◀ ref" if a.peak_index == int(ref_i) else ""
                rows.append([f"{a.peak_index+1}{mark}", f"{a.observed:.5f}",
                             f"{a.delta_obs_mau:+.2f}", a.candidate.formula,
                             f"{a.delta_theo_mau:+.2f}", f"{a.err_mau:+.2f}",
                             f"{a.candidate.plausibility:.4f}"])
            header = ["#", "observed", "Δobs (mau)", "candidate",
                      "Δtheo (mau)", "err (mau)", "plaus"]
            objs = [pn.pane.Markdown(
                f"Reference: peak {int(ref_i)+1} = **{ref_f}** · differences "
                "from it; scale assumed 1 (no absolute calibration)."
                + ("  ·  *natural abundances ignored*" if labeled else ""))]
            objs.append(_html_table(header, rows, width=760, right_from=4))
            id_out.objects = objs
            H["ident_labels"] = [(a.observed, a.candidate.formula) for a in res_id]
            H["ident_text"] = (
                f"Species ID — reference peak {int(ref_i)+1} = {ref_f} "
                f"(differences, scale 1"
                + ("; natural abundances ignored" if labeled else "") + ")\n\n"
                + _text_table(header, rows))
            H["ident_caption"] = (
                f"HMR species ID (reference {ref_f}) — "
                f"{_os.path.basename(H['path'])}")

        else:   # blind best-guess
            try:
                sols = identify_blind(peaks_mz, elems, polarity=pol,
                                      max_atoms=maxa, nominal_mass=H.get("nominal"),
                                      ignore_abundance=labeled, top_n=5)
            except Exception as exc:
                id_out.objects = [pn.pane.Markdown(f"❌ Identify failed: `{exc}`")]
                return
            if not sols:
                id_out.objects = [pn.pane.Markdown(
                    "No candidates — check the selected elements / polarity.")]
                return
            gap = (sols[1].rms_mau - sols[0].rms_mau) if len(sols) > 1 else float("inf")
            conf = ("clear winner" if gap >= 0.75 else
                    "weak — assignments nearly tie" if gap < 0.3 else "moderate")
            rows = []
            for rank, s in enumerate(sols):
                rows.append([f"{rank+1}", f"{s.rms_mau:.2f}",
                             " | ".join(a.formula for a in s.assignment),
                             f"{s.mean_plausibility:.4f}"])
            header = ["rank", "RMS (mau)", "assignment (peak 1 → N)", "mean plaus"]
            objs = [pn.pane.Markdown(
                f"Best guess from differences · confidence: **{conf}** "
                f"(gap to next = {gap:.2f} mau)."
                + ("  ·  *natural abundances ignored*" if labeled else "")
                + "  <span style='font-size:11px;opacity:0.7'>Fewer elements ⇒ "
                "fewer candidates ⇒ a more trustworthy gap.</span>")]
            objs.append(_html_table(header, rows, width=820, right_from=1))
            id_out.objects = objs
            best = sols[0]
            H["ident_labels"] = [(p, a.formula)
                                 for p, a in zip(peaks_mz, best.assignment)]
            H["ident_text"] = (
                f"Species ID — blind best guess (confidence: {conf}, "
                f"gap {gap:.2f} mau"
                + ("; natural abundances ignored" if labeled else "") + ")\n\n"
                + _text_table(header, rows))
            H["ident_caption"] = (
                f"HMR species ID (blind, {conf}) — "
                f"{_os.path.basename(H['path'])}")
        _render_spectrum()
    id_btn.on_click(_identify)

    # ── journal / export ───────────────────────────────────────────────────────
    jr_spec = pn.widgets.Button(name="➕ Spectrum → journal", width=180)
    jr_resid = pn.widgets.Button(name="➕ Residuals → journal", width=180)
    jr_peaks = pn.widgets.Button(name="➕ Peak table → journal", width=180)
    jr_id = pn.widgets.Button(name="➕ Species ID → journal", width=180)
    jr_save = pn.widgets.Button(name="💾 Save spectrum PNG", button_type="default",
                                width=180)
    jr_status = pn.pane.Markdown("")

    def _fig_png(fig, dpi=150):
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight")
        plt.close(fig)
        return buf.getvalue()

    def _jr_add_spectrum(event=None):
        fig = _build_spectrum_fig()
        if fig is None:
            jr_status.object = "⚠️ Load a scan first."
            return
        cap = f"HMR spectrum ({x_axis.value}) — {_os.path.basename(H['path'])}"
        if H["fit"] is not None:
            cap += f" · {H['fit'].n_peaks}-peak fit"
        JOURNAL.append({"png": _fig_png(fig), "caption": cap,
                        "added": datetime.datetime.now().strftime("%Y-%m-%d %H:%M")})
        _update_journal_status()
        jr_status.object = f"✅ Spectrum added — {len(JOURNAL)} item(s)."
    jr_spec.on_click(_jr_add_spectrum)

    def _jr_add_resid(event=None):
        fig = _build_residuals_fig()
        if fig is None:
            jr_status.object = "⚠️ Run a deconvolution first."
            return
        cap = f"HMR deconvolution residuals — {_os.path.basename(H['path'])}"
        JOURNAL.append({"png": _fig_png(fig), "caption": cap,
                        "added": datetime.datetime.now().strftime("%Y-%m-%d %H:%M")})
        _update_journal_status()
        jr_status.object = f"✅ Residuals added — {len(JOURNAL)} item(s)."
    jr_resid.on_click(_jr_add_resid)

    def _jr_add_peaks(event=None):
        res = H["fit"]
        if res is None or res.n_peaks == 0:
            jr_status.object = "⚠️ Run a deconvolution first."
            return
        ax_key = H["fit_axis"]
        anchor = res.peaks[0].mu
        rows = [[f"{i+1}", f"{p.mu:.5f}", f"{p.a:.0f}", f"±{p.a_se:.0f}",
                 f"{(p.mu-anchor)*1000:+.2f}",
                 f"{(p.a/p.a_se):.1f}" if p.a_se > 0 else "—"]
                for i, p in enumerate(res.peaks)]
        header = ["#", f"centre({ax_key})", "amp", "amp±",
                  f"Δ#1(m{ax_key})", "a/σ"]
        txt = (f"HMR deconvolution — w={res.w*1000:.3f} σ={res.sigma*1000:.3f} "
               f"m{ax_key}, χ²={res.chi2:.1f}, BIC={res.bic:.1f}\n\n"
               + _text_table(header, rows))
        _journal_add_text(txt, f"HMR peak table — {_os.path.basename(H['path'])}")
        jr_status.object = f"✅ Peak table added — {len(JOURNAL)} item(s)."
    jr_peaks.on_click(_jr_add_peaks)

    def _jr_add_id(event=None):
        if not H.get("ident_text"):
            jr_status.object = "⚠️ Run an identification first."
            return
        _journal_add_text(H["ident_text"], H["ident_caption"])
        jr_status.object = f"✅ Species ID added — {len(JOURNAL)} item(s)."
    jr_id.on_click(_jr_add_id)

    def _jr_save_png(event=None):
        fig = _build_spectrum_fig()
        if fig is None:
            jr_status.object = "⚠️ Load a scan first."
            return
        folder = _os.path.expanduser(journal_dir.value or "~")
        stem = _os.path.splitext(_os.path.basename(H["path"]))[0]
        out = _os.path.join(folder, f"{stem}_spectrum.png")
        try:
            fig.savefig(out, dpi=200, bbox_inches="tight"); plt.close(fig)
        except Exception as exc:
            jr_status.object = f"❌ Save failed: `{exc}`"
            return
        jr_status.object = f"✅ Saved `{out}`"
    jr_save.on_click(_jr_save_png)

    _toggle_ref_row()       # set initial reference-row visibility
    _scan_folder()          # populate file list on first build

    note = pn.pane.Markdown(
        "<span style='font-size:11px;opacity:0.75'>The amu axis is the "
        "instrument's nominal scale, not exact mass, so identification uses "
        "<i>relative</i> peak separations only (scale assumed 1, no absolute "
        "calibration). <b>Reference peak</b>: you tag one peak you know and the "
        "rest are read off by mass difference from it. <b>Blind</b>: best guess "
        "from the separation pattern with nothing known — trust it only when the "
        "RMS gap to the next assignment is wide, which needs a tight element "
        "set. The labeled-sample box drops natural abundances from the ranking. "
        "MRP uses the 10–90% edge width.</span>")

    return pn.Column(
        pn.pane.Markdown("### HMR — high mass resolution"),
        pn.Row(hmr_dir, hmr_scan_btn),
        pn.Row(hmr_file, hmr_load_btn),
        hmr_info,
        pn.layout.Divider(),
        pn.pane.Markdown("**Spectrum**"),
        pn.Row(x_axis, log_y),
        spectrum_box,
        pn.layout.Divider(),
        pn.pane.Markdown("**IRF & resolving power**"),
        pn.Row(edge_side, m_ref, mrp_btn),
        mrp_out,
        pn.layout.Divider(),
        pn.pane.Markdown("**Deconvolution**"),
        pn.Row(dc_method, dc_run),
        pn.Row(dc_nmin, dc_nmax, dc_grid, dc_stab, dc_minsnr),
        dc_status,
        dc_table,
        resid_box,
        pn.layout.Divider(),
        pn.pane.Markdown("**Species identification** "
                         "<span style='font-size:11px;opacity:0.7'>"
                         "(relative separation)</span>"),
        pn.pane.Markdown("<span style='font-size:11px;opacity:0.7'>Click "
                         "elements expected in the sample:</span>"),
        periodic_table,
        pn.Row(sel_label, id_clear),
        pn.Row(id_pol, id_maxatoms, id_labeled),
        id_mode,
        id_ref_row,
        pn.Row(id_btn),
        id_out,
        pn.layout.Divider(),
        pn.pane.Markdown("**Journal & export**"),
        pn.Row(jr_spec, jr_resid, jr_peaks, jr_id),
        pn.Row(jr_save, jr_status),
        note,
    )


def _cluster_summary_text(result, k):
    """Compact text summary of a ClusterResult at a chosen k, suitable for
    the journal. Mirrors the on-screen cluster table in plain text."""
    import numpy as _np
    lines = []
    lines.append(f"Clustering — {result['method']}, "
                 f"feature space: {result['feature_space']}")
    lines.append(f"channels: {', '.join(result['feature_labels'])}")
    lines.append(f"sensible k = {result['sensible_k']}   (showing k = {k})")
    if result['method'] == 'hierarchical' and 'cophenetic_corr' in result:
        lines.append(f"cophenetic correlation = {result['cophenetic_corr']:.3f}")
    recs = ", ".join(
        f"k={r['k']} [{', '.join(m.replace('_', ' ') for m in r['methods'])}]"
        for r in result['unique_k_recommendations'])
    lines.append(f"recommendations: {recs}")
    lines.append("")

    sizes = result['cluster_sizes_by_k'][k]
    cents = result['centroids_by_k_counts'][k]
    feats = result['feature_labels']
    total = int(sum(sizes))
    header = f"{'#':>2}  {'pixels':>9}  {'%':>6}  " + \
             "  ".join(f"{f:>12}" for f in feats)
    lines.append(header)
    lines.append("-" * len(header))
    for i in range(k):
        size = int(sizes[i]) if i < len(sizes) else 0
        pct = 100.0 * size / total if total else 0.0
        row_c = cents[i] if i < cents.shape[0] else _np.zeros(len(feats))
        cols = "  ".join(f"{v:>12.1f}" for v in row_c)
        lines.append(f"{i + 1:>2}  {size:>9,}  {pct:>5.1f}%  {cols}")
    return "\n".join(lines)


def _text_table(header, rows):
    """Render a table as aligned monospace text (for journal text entries)."""
    cols = list(zip(*([header] + rows))) if rows else [[h] for h in header]
    widths = [max(len(str(c)) for c in col) for col in cols]
    def fmt(r):
        return "  ".join(str(v).ljust(widths[i]) for i, v in enumerate(r))
    line = "  ".join("-" * w for w in widths)
    return "\n".join([fmt(header), line] + [fmt(r) for r in rows])


def _html_table(header, rows, width=620, right_from=1):
    """A theme-agnostic table that stays readable in light or dark mode.

    Inherits the surrounding text colour (so no dark-on-dark) and uses a faint
    translucent stripe instead of an opaque alternating background. Numeric
    columns (index >= ``right_from``) are right-aligned with tabular figures.
    Replaces pn.pane.DataFrame, whose default striping renders unreadable on the
    dark theme.
    """
    def align(c):
        return "right" if c >= right_from else "left"
    th = "".join(
        f"<th style='text-align:{align(c)};padding:5px 12px;"
        f"border-bottom:1px solid rgba(128,128,128,0.5);font-weight:600;'>{h}</th>"
        for c, h in enumerate(header))
    trs = []
    for i, r in enumerate(rows):
        bg = "rgba(128,128,128,0.12)" if i % 2 else "transparent"
        tds = "".join(
            f"<td style='text-align:{align(c)};padding:4px 12px;'>{v}</td>"
            for c, v in enumerate(r))
        trs.append(f"<tr style='background:{bg};'>{tds}</tr>")
    html = (
        f"<table style='border-collapse:collapse;color:inherit;font-size:13px;"
        f"font-variant-numeric:tabular-nums;min-width:{int(width*0.6)}px;'>"
        f"<thead><tr>{th}</tr></thead><tbody>{''.join(trs)}</tbody></table>")
    return pn.pane.HTML(html, width=width)


def _enlarge_for_display(fig, per_w=3.1, per_h=3.1, cap_w=34, cap_h=44):
    """Scale a matplotlib figure up so multi-panel grids render legibly.

    Infers the subplot grid from the first axes' gridspec geometry and sizes the
    figure at roughly `per_w`×`per_h` inches per panel, capped. Returns the
    figure's pixel dimensions at `dpi` so the caller can size the display pane.
    """
    try:
        nrows, ncols = fig.axes[0].get_gridspec().get_geometry()
    except Exception:
        nrows, ncols = 1, max(1, len(fig.axes))
    w_in = min(cap_w, max(9.0, ncols * per_w))
    h_in = min(cap_h, max(3.2, nrows * per_h))
    fig.set_size_inches(w_in, h_in)
    try:
        fig.tight_layout()
    except Exception:
        pass
    return w_in, h_in


def analysis_view(img):
    if img is None:
        return pn.pane.Markdown(
            "### Analysis\nLoad a `.im` file to run clustering and distributions.")

    masses = img.masses
    non_se = [m for m in masses if "SE" not in m.upper()] or masses

    drift_hint = ("" if img.corrected is not None else
                  "<span style='color:#e0a030'>Tip: run drift correction on "
                  "the Imaging tab first — clustering and GMM fits use the "
                  "drift-corrected stack.</span>")

    def _fig_to_png(fig, dpi=150):
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight",
                    facecolor=fig.get_facecolor())
        plt.close(fig)
        return buf.getvalue()

    # ── Clustering subsection ────────────────────────────────────────────────
    cl_method = pn.widgets.RadioButtonGroup(
        name="Method", options=["kmeans", "hierarchical"],
        value="kmeans", button_type="primary")
    cl_kmax = pn.widgets.IntInput(name="k max", value=10, start=2, end=20, width=90)
    cl_feature = pn.widgets.Select(
        name="Feature space",
        options=["log_zscored", "log_robustz", "log", "raw"],
        value="log_zscored", width=150)
    cl_channels = pn.widgets.MultiSelect(
        name="Channels (none = auto, SE excluded)",
        options=list(masses), value=list(non_se), size=4, width=220)
    cl_include_se = pn.widgets.Checkbox(name="Include SE (auto mode only)",
                                        value=False)
    cl_mincounts = pn.widgets.IntInput(name="Min counts (0 = off)",
                                        value=0, start=0, width=150)
    cl_maskchan = pn.widgets.Select(name="Mask channel",
                                    options=["auto"] + list(masses),
                                    value="auto", width=150)
    cl_subsample = pn.widgets.IntInput(name="Subsample (hier.)",
                                       value=5000, start=200, width=140)
    cl_linkage = pn.widgets.Select(
        name="Linkage (hier.)",
        options=["ward", "complete", "average", "single", "centroid"],
        value="ward", width=140)

    cl_run_btn = pn.widgets.Button(name="Run clustering", button_type="primary",
                                   icon="player-play", width=150)
    cl_kpick = pn.widgets.Select(name="Show k", options=[], width=90,
                                 disabled=True)
    cl_display = pn.widgets.Select(
        name="Display",
        options=["Cluster labels", "Cluster grid", "Overlay on channel",
                 "Metric sweep", "Dendrogram (hier.)"],
        value="Cluster labels", width=180)
    cl_overlay_chan = pn.widgets.Select(name="Overlay channel",
                                        options=list(masses),
                                        value=non_se[0], width=150)
    cl_add_fig_btn = pn.widgets.Button(name="➕ Add figure to journal",
                                       button_type="default", width=200)
    cl_add_txt_btn = pn.widgets.Button(name="➕ Add summary to journal",
                                       button_type="default", width=210)
    cl_status = pn.pane.Markdown("")
    cl_plot_box = pn.Column(pn.pane.Markdown(
        "<span style='opacity:0.6'>No clustering yet — set options and click "
        "**Run clustering**.</span>"))

    _cl = {"result": None, "png": None}

    def _cl_render():
        result = _cl["result"]
        if result is None:
            return
        try:
            k = int(cl_kpick.value)
        except (TypeError, ValueError):
            k = result["sensible_k"]
        disp = cl_display.value
        try:
            if disp == "Cluster labels":
                fig = plot_cluster_labels(img, result, k=k, show=False)
            elif disp == "Cluster grid":
                fig = plot_cluster_grid(img, result, show=False)
            elif disp == "Overlay on channel":
                fig = plot_overlay(img, result, k=k, base="channel",
                                   channel=cl_overlay_chan.value, show=False)
            elif disp == "Metric sweep":
                fig = plot_metric_sweep(result, show=False)
            else:   # Dendrogram
                if result["method"] != "hierarchical":
                    cl_plot_box.objects = [pn.pane.Alert(
                        "Dendrogram is only available for hierarchical "
                        "clustering.", alert_type="warning")]
                    return
                fig = plot_dendrogram(result, k_marks=[k], show=False)
        except Exception as exc:
            cl_plot_box.objects = [pn.pane.Alert(
                f"Render failed: `{exc}`", alert_type="danger")]
            return
        _cl["png"] = _fig_to_png(fig)
        cl_plot_box.objects = [pn.pane.PNG(_cl["png"], width=720)]

    def _cl_run(event):
        ch_list = list(cl_channels.value)
        channels_arg = ch_list if ch_list else None
        try:
            result = cluster_pixels(
                img, method=cl_method.value, k_max=int(cl_kmax.value),
                channels=channels_arg, include_se=bool(cl_include_se.value),
                feature_space=cl_feature.value,
                min_counts=(int(cl_mincounts.value) or None),
                mask_channel=(None if cl_maskchan.value == "auto"
                              else cl_maskchan.value),
                subsample_size=int(cl_subsample.value),
                linkage_method=cl_linkage.value, verbose=False)
        except Exception as exc:
            cl_status.object = f"❌ Clustering failed: `{exc}`"
            return
        _cl["result"] = result
        ks = sorted(int(x) for x in result["labels_by_k"])
        cl_kpick.options = [str(x) for x in ks]
        cl_kpick.value = str(result["sensible_k"])
        cl_kpick.disabled = False
        recs = ", ".join(
            f"k={r['k']} ({', '.join(m.replace('_', ' ') for m in r['methods'])})"
            for r in result["unique_k_recommendations"])
        cl_status.object = (f"✅ Done — sensible k = **{result['sensible_k']}**. "
                            f"Recommendations: {recs}.")
        _cl_render()
    cl_run_btn.on_click(_cl_run)

    cl_kpick.param.watch(lambda e: _cl_render(), "value")
    cl_display.param.watch(lambda e: _cl_render(), "value")
    cl_overlay_chan.param.watch(lambda e: _cl_render(), "value")

    def _cl_add_fig(event):
        if _cl["png"] is None:
            cl_status.object = "Nothing to add yet — run clustering first."
            return
        k = cl_kpick.value
        res = _cl["result"]
        cap = (f"{res['method']} clustering, {cl_display.value.lower()}, k={k}, "
               f"{res['feature_space']}  ·  {os.path.basename(img.path)}")
        JOURNAL.append({"png": _cl["png"], "caption": cap,
                        "added": datetime.datetime.now().strftime("%Y-%m-%d %H:%M")})
        _update_journal_status()
        cl_status.object = f"✅ Figure added — {len(JOURNAL)} item(s) in journal."
    cl_add_fig_btn.on_click(_cl_add_fig)

    def _cl_add_txt(event):
        res = _cl["result"]
        if res is None:
            cl_status.object = "Nothing to add yet — run clustering first."
            return
        try:
            k = int(cl_kpick.value)
        except (TypeError, ValueError):
            k = res["sensible_k"]
        text = _cluster_summary_text(res, k)
        cap = (f"{res['method']} cluster summary, k={k}  ·  "
               f"{os.path.basename(img.path)}")
        _journal_add_text(text, cap)
        cl_status.object = f"✅ Summary added — {len(JOURNAL)} item(s) in journal."
    cl_add_txt_btn.on_click(_cl_add_txt)

    clustering_section = pn.Column(
        pn.pane.Markdown("**Clustering**  <span style='font-size:11px;"
                         "opacity:0.7'>(k-means / hierarchical, metric sweep)</span>"),
        pn.Row(cl_method, cl_kmax, cl_feature),
        pn.Row(cl_channels, pn.Column(cl_include_se, cl_mincounts, cl_maskchan)),
        pn.Row(cl_subsample, cl_linkage),
        pn.Row(cl_run_btn, cl_kpick, cl_display, cl_overlay_chan),
        pn.Row(cl_add_fig_btn, cl_add_txt_btn),
        cl_status,
        cl_plot_box,
    )

    # ── Distributions (per-channel GMM histograms) subsection ────────────────
    h_channel = pn.widgets.Select(name="Channel",
                                  options=["all channels"] + list(masses),
                                  value=non_se[0], width=170)
    h_kmax = pn.widgets.IntInput(name="k max", value=6, start=1, end=10, width=90)
    h_nbins = pn.widgets.IntInput(name="Bins", value=80, start=20, end=300, width=90)
    h_dropzeros = pn.widgets.Checkbox(name="Drop zero-count pixels", value=True)
    h_jitter = pn.widgets.Checkbox(name="Jitter integer counts", value=True)
    h_tailw = pn.widgets.FloatInput(name="Tail-warning weight",
                                    value=0.15, start=0.0, step=0.05, width=160)
    h_run_btn = pn.widgets.Button(name="Run distributions",
                                  button_type="primary", icon="chart-histogram",
                                  width=170)
    h_add_fig_btn = pn.widgets.Button(name="➕ Add figure to journal",
                                      button_type="default", width=200)
    h_add_txt_btn = pn.widgets.Button(name="➕ Add summary to journal",
                                      button_type="default", width=210)
    h_status = pn.pane.Markdown("")
    h_plot_box = pn.Column(pn.pane.Markdown(
        "<span style='opacity:0.6'>No fit yet — set options and click "
        "**Run distributions**.</span>"))

    _h = {"png": None, "text": None}

    def _h_run(event):
        ch_arg = None if h_channel.value == "all channels" else h_channel.value
        plt.close("all")   # clean slate so plt.gcf() is unambiguous below
        buf_txt = io.StringIO()
        try:
            with contextlib.redirect_stdout(buf_txt):
                plot_histograms(
                    img, channel=ch_arg, k_max=int(h_kmax.value),
                    n_bins=int(h_nbins.value),
                    drop_zeros=bool(h_dropzeros.value),
                    jitter=bool(h_jitter.value),
                    tail_weight_threshold=float(h_tailw.value),
                    show=True, verbose=True)
        except Exception as exc:
            h_status.object = f"❌ Distribution fit failed: `{exc}`"
            return
        fig = plt.gcf()
        w_in, h_in = _enlarge_for_display(fig)        # bigger panels
        dpi = 150
        px_w, px_h = int(w_in * dpi), int(h_in * dpi)
        _h["png"] = _fig_to_png(fig, dpi=dpi)         # also closes the figure
        _h["text"] = buf_txt.getvalue().rstrip()
        # Show at natural (large) size inside a scroll box — horizontal scroll
        # for the wide single-channel row, vertical for the tall all-channels
        # stack. Capped height keeps the page navigable.
        h_plot_box.height = min(px_h + 24, 720)
        h_plot_box.scroll = True
        h_plot_box.objects = [pn.pane.PNG(_h["png"], width=px_w)]
        lbl = "all channels" if ch_arg is None else ch_arg
        h_status.object = (f"✅ GMM fit complete for {lbl}. "
                           f"<span style='opacity:0.65;font-size:11px'>"
                           f"(scroll the panel; fewer **k max** = larger panels)</span>")
    h_run_btn.on_click(_h_run)

    def _h_add_fig(event):
        if _h["png"] is None:
            h_status.object = "Nothing to add yet — run a fit first."
            return
        lbl = "all channels" if h_channel.value == "all channels" else h_channel.value
        cap = (f"GMM histograms — {lbl}, k≤{int(h_kmax.value)}  ·  "
               f"{os.path.basename(img.path)}")
        JOURNAL.append({"png": _h["png"], "caption": cap,
                        "added": datetime.datetime.now().strftime("%Y-%m-%d %H:%M")})
        _update_journal_status()
        h_status.object = f"✅ Figure added — {len(JOURNAL)} item(s) in journal."
    h_add_fig_btn.on_click(_h_add_fig)

    def _h_add_txt(event):
        if not _h["text"]:
            h_status.object = "No summary captured — run a fit first."
            return
        lbl = "all channels" if h_channel.value == "all channels" else h_channel.value
        cap = f"GMM threshold summary — {lbl}  ·  {os.path.basename(img.path)}"
        _journal_add_text(_h["text"], cap)
        h_status.object = f"✅ Summary added — {len(JOURNAL)} item(s) in journal."
    h_add_txt_btn.on_click(_h_add_txt)

    distributions_section = pn.Column(
        pn.pane.Markdown("**Distributions**  <span style='font-size:11px;"
                         "opacity:0.7'>(per-channel GMM histograms + BIC/AIC + "
                         "crossing thresholds)</span>"),
        pn.Row(h_channel, h_kmax, h_nbins),
        pn.Row(h_dropzeros, h_jitter, h_tailw),
        pn.Row(h_run_btn, h_add_fig_btn, h_add_txt_btn),
        h_status,
        h_plot_box,
    )

    # ── ROI Manager subsection ───────────────────────────────────────────────
    # Two ways to make ROIs: (1) from a clustering run's per-cluster masks, or
    # (2) by edge/threshold detection on a chosen channel. ROIs then get
    # per-channel pooled-count statistics (and optional pooled isotope ratios).
    _roi = {"rois": {}, "pending": {}, "stats_png": None, "stats_text": None}

    def _summed_channel(channel, corrected=True):
        """Summed-over-planes (H, W) counts for one channel.

        corrected=True follows get_channel (the drift-corrected float stack if
        drift has been run, else the raw stack). corrected=False forces the raw
        integer-count stack regardless of drift state — use for Poisson error
        propagation, where you need true detected-ion counts."""
        if corrected:
            stack = np.asarray(img.get_channel(channel))
            return (stack.sum(axis=0) if stack.ndim == 3 else stack).astype(float)
        ch = list(img.masses).index(channel)
        return img.data[:, ch, :, :].sum(axis=0).astype(float)

    def _masks_from_clusters(result, k):
        """{cluster_id (1-based int): (H, W) bool mask} from the library helper.
        IDs and pixel counts match the clustering summary table; pixels excluded
        by min_counts / mask_channel are NaN in labels_by_k and fall out of
        every mask."""
        return extract_cluster_masks(result, k=k)

    def _refresh_roi_list():
        rois = _roi["rois"]
        # Attach the current masks to the image object so the stack carries its
        # ROI set (used by depth profiles here, and available to any downstream
        # consumer holding the same MimsImage). Not written to .npz sessions —
        # that would need a core save_session change.
        img.rois = {name: info["mask"] for name, info in rois.items()}
        # keep the "Show on map" dropdown in step with the live ROI set
        roi_show.options = ["all"] + list(rois.keys())
        if roi_show.value not in roi_show.options:
            roi_show.value = "all"
        if not rois:
            roi_list.object = "<span style='opacity:0.6'>No ROIs yet.</span>"
            return
        lines = [f"**ROIs — {len(rois)}:**"]
        for i, (name, info) in enumerate(rois.items(), 1):
            lines.append(f"- **#{i}** `{name}` · {int(info['mask'].sum()):,} px "
                         f"<span style='opacity:0.6'>({info['source']})</span>")
        roi_list.object = "  \n".join(lines)

    # -- (1) From clusters ----------------------------------------------------
    roi_load_btn = pn.widgets.Button(name="Load clusters from last run",
                                     button_type="default", width=220)
    roi_cluster_pick = pn.widgets.MultiSelect(name="Clusters to add",
                                              options=[], size=5, width=240)
    roi_add_clusters_btn = pn.widgets.Button(name="➕ Add selected as ROIs",
                                             button_type="primary", width=200)

    def _roi_load_clusters(event):
        result = _cl["result"]
        if result is None:
            roi_status.object = "⚠️ Run clustering first (Clustering section above)."
            return
        try:
            k = int(cl_kpick.value)
        except (TypeError, ValueError):
            k = result["sensible_k"]
        try:
            masks = _masks_from_clusters(result, k)
        except Exception as exc:
            roi_status.object = f"❌ Couldn't build cluster masks: `{exc}`"
            return
        _roi["pending"] = {f"k{k}_c{cid}": m for cid, m in sorted(masks.items())}
        roi_cluster_pick.options = [
            f"{name}  ({int(m.sum()):,} px)"
            for name, m in _roi["pending"].items()]
        roi_cluster_pick.value = list(roi_cluster_pick.options)
        roi_status.object = (f"Loaded {len(masks)} cluster(s) at k={k}. "
                             f"Pick which to add, then **Add selected as ROIs**.")
    roi_load_btn.on_click(_roi_load_clusters)

    def _roi_add_clusters(event):
        if not _roi["pending"]:
            roi_status.object = "Nothing loaded — click **Load clusters** first."
            return
        label_to_name = {f"{name}  ({int(m.sum()):,} px)": name
                         for name, m in _roi["pending"].items()}
        added = 0
        for lbl in roi_cluster_pick.value:
            name = label_to_name.get(lbl)
            if name:
                _roi["rois"][name] = {"mask": _roi["pending"][name],
                                      "source": "cluster"}
                added += 1
        _refresh_roi_list()
        roi_status.object = f"✅ Added {added} ROI(s) — {len(_roi['rois'])} total."
    roi_add_clusters_btn.on_click(_roi_add_clusters)

    # -- (1b) From the Draw ROIs tab -----------------------------------------
    # Drawn shapes are staged onto img._pending_drawn by the Draw tab's
    # "Send to Analysis" button; here they enter the ROI set as a third source.
    roi_load_drawn_btn = pn.widgets.Button(
        name="⬇ Load drawn ROIs from Draw tab", button_type="primary", width=260)

    def _roi_load_drawn(event):
        pending = getattr(img, "_pending_drawn", None)
        if not pending:
            roi_status.object = ("No drawn ROIs staged — go to the **Draw ROIs** "
                                 "tab, draw shapes, **Rasterise**, then "
                                 "**Send to Analysis**.")
            return
        added = 0
        for name, m in pending.items():
            key = name if name not in _roi["rois"] else f"{name}_{len(_roi['rois'])}"
            _roi["rois"][key] = {"mask": np.asarray(m, bool), "source": "drawn"}
            added += 1
        _refresh_roi_list()
        roi_status.object = (f"✅ Loaded {added} drawn ROI(s) — "
                             f"{len(_roi['rois'])} total. Set channels and click "
                             f"**Compute ROI stats** below.")
    roi_load_drawn_btn.on_click(_roi_load_drawn)

    # -- (2) Edge / threshold detection --------------------------------------
    roi_edge_chan = pn.widgets.Select(name="Detect on channel",
                                      options=list(masses), value=non_se[0],
                                      width=160)
    roi_edge_method = pn.widgets.Select(
        name="Method",
        options=["Threshold (Otsu)", "Manual percentile", "Edge (Canny) + fill"],
        value="Threshold (Otsu)", width=190)
    roi_edge_sigma = pn.widgets.FloatInput(name="Smoothing σ", value=1.5,
                                           start=0.0, step=0.5, width=120)
    roi_edge_pct = pn.widgets.FloatInput(name="Percentile (manual)", value=90.0,
                                         start=0.0, end=100.0, step=1.0, width=160)
    roi_edge_minarea = pn.widgets.IntInput(name="Min area (px)", value=25,
                                           start=1, width=130)
    roi_edge_merge = pn.widgets.Checkbox(name="Merge into one ROI", value=False)
    roi_detect_btn = pn.widgets.Button(name="Detect regions",
                                       button_type="primary", icon="vector",
                                       width=160)
    roi_overlay_box = pn.Column(pn.pane.Markdown(
        "<span style='opacity:0.6'>No detection yet.</span>"))

    def _detect_regions():
        from scipy import ndimage as ndi
        data = _summed_channel(roi_edge_chan.value)
        norm = data / (data.max() or 1.0)
        sigma = float(roi_edge_sigma.value)
        method = roi_edge_method.value
        if method == "Edge (Canny) + fill":
            from skimage.feature import canny
            binary = ndi.binary_fill_holes(canny(norm, sigma=sigma))
        else:
            from skimage.filters import gaussian
            sm = gaussian(norm, sigma=sigma) if sigma > 0 else norm
            if method == "Threshold (Otsu)":
                from skimage.filters import threshold_otsu
                try:
                    thr = threshold_otsu(sm)
                except Exception:
                    thr = float(sm.mean())
            else:                                   # Manual percentile
                thr = float(np.percentile(sm, float(roi_edge_pct.value)))
            binary = ndi.binary_fill_holes(sm > thr)
        if bool(roi_edge_merge.value):
            return ({"edge_all": binary} if binary.any() else {}), data
        lab, n = ndi.label(binary)
        regions = [(lab == i) for i in range(1, n + 1)]
        regions = [m for m in regions if m.sum() >= int(roi_edge_minarea.value)]
        regions.sort(key=lambda m: -m.sum())
        return {f"r{j}": m for j, m in enumerate(regions[:50], 1)}, data

    def _roi_overlay_fig(base_channel, extra_masks=None, only=None,
                         numbered=True):
        data = _summed_channel(base_channel)
        vmax = float(np.percentile(data, 99.5)) or 1.0
        fig, ax = plt.subplots(figsize=(5.8, 5.8))
        ax.imshow(data, cmap="gray", vmin=0, vmax=vmax, interpolation="nearest")
        cmap = plt.get_cmap("tab10")
        items = list((extra_masks or _roi["rois"]).items())
        shown = 0
        for i, (name, info) in enumerate(items):
            if only is not None and name != only:
                continue
            m = info["mask"] if isinstance(info, dict) else info
            col = cmap(i % 10)
            ax.contour(m.astype(float), levels=[0.5], colors=[col],
                       linewidths=1.4)
            if numbered:
                ys_, xs_ = np.where(m)
                if len(xs_):
                    # the number is the link to the "#" column in the stats table
                    ax.text(xs_.mean(), ys_.mean(), str(i + 1),
                            color="white", fontsize=10, fontweight="bold",
                            ha="center", va="center",
                            bbox=dict(boxstyle="circle,pad=0.3", fc=col,
                                      ec="white", lw=0.8, alpha=0.95))
            shown += 1
        if only is not None:
            names = [n for n, _ in items]
            idx = names.index(only) + 1 if only in names else 0
            ax.set_title(f"ROI #{idx} — {only}  on {base_channel}", fontsize=10)
        else:
            ax.set_title(f"{shown} region(s) on {base_channel}  "
                         f"(numbers = table #)", fontsize=10)
        ax.set_xlabel("px"); ax.set_ylabel("px")
        fig.tight_layout()
        return fig

    def _roi_detect(event):
        try:
            found, _ = _detect_regions()
        except Exception as exc:
            roi_status.object = f"❌ Detection failed: `{exc}`"
            return
        if not found:
            roi_overlay_box.objects = [pn.pane.Markdown(
                "<span style='opacity:0.6'>Nothing detected — lower the "
                "threshold / min area.</span>")]
            roi_status.object = "No regions found."
            return
        # add directly to the ROI set, tagging the source
        for name, m in found.items():
            _roi["rois"][name] = {"mask": m, "source": "edge"}
        fig = _roi_overlay_fig(roi_edge_chan.value)
        png = _fig_to_png(fig)
        roi_overlay_box.objects = [pn.pane.PNG(png, width=560)]
        _refresh_roi_list()
        roi_status.object = (f"✅ Detected & added {len(found)} region(s) — "
                             f"{len(_roi['rois'])} ROI(s) total.")
    roi_detect_btn.on_click(_roi_detect)

    # -- ROI stats ------------------------------------------------------------
    roi_stat_chans = pn.widgets.MultiSelect(
        name="Stat channels", options=list(masses), value=list(non_se),
        size=4, width=200)
    roi_counts_src = pn.widgets.Select(
        name="Counts source",
        options=["Drift-corrected (if run)", "Raw counts (integer)"],
        value="Drift-corrected (if run)", width=210)
    roi_ratio_on = pn.widgets.Checkbox(name="Add pooled ratio", value=False)
    roi_ratio_num = pn.widgets.Select(name="Ratio numerator",
                                      options=list(masses), value=non_se[0],
                                      width=150)
    roi_ratio_den = pn.widgets.Select(
        name="Ratio denominator", options=list(masses),
        value=(non_se[1] if len(non_se) > 1 else non_se[0]), width=150)
    roi_compute_btn = pn.widgets.Button(name="Compute ROI stats",
                                        button_type="primary", icon="table",
                                        width=180)
    roi_remove_btn = pn.widgets.Button(name="Remove last ROI",
                                       button_type="default", width=150)
    roi_clear_btn = pn.widgets.Button(name="Clear ROIs",
                                      button_type="default", width=130)
    roi_add_fig_btn = pn.widgets.Button(name="➕ Add table to journal",
                                        button_type="default", width=200)
    roi_add_txt_btn = pn.widgets.Button(name="➕ Add summary to journal",
                                        button_type="default", width=210)
    roi_table_box = pn.Column(pn.pane.Markdown(
        "<span style='opacity:0.6'>No stats yet — add ROIs and click "
        "**Compute ROI stats**.</span>"))

    # -- ROI map (numbered overlay; numbers correspond to the table "#" col) --
    roi_map_chan = pn.widgets.Select(name="Map channel", options=list(masses),
                                     value=non_se[0], width=160)
    roi_show = pn.widgets.Select(name="Show on map", options=["all"],
                                 value="all", width=200)
    roi_map_box = pn.Column(pn.pane.Markdown(
        "<span style='opacity:0.6'>ROI map appears here after **Compute ROI "
        "stats** — each region is numbered to match the table.</span>"))

    def _render_map(*_):
        if not _roi["rois"]:
            roi_map_box.objects = [pn.pane.Markdown(
                "<span style='opacity:0.6'>No ROIs to map.</span>")]
            return
        only = None if roi_show.value == "all" else roi_show.value
        try:
            fig = _roi_overlay_fig(roi_map_chan.value, only=only)
        except Exception as exc:
            roi_status.object = f"❌ Map failed: `{exc}`"
            return
        roi_map_box.objects = [pn.pane.PNG(_fig_to_png(fig), width=560)]
    roi_show.param.watch(_render_map, "value")
    roi_map_chan.param.watch(_render_map, "value")

    def _build_stats_rows(channels, ratio, corrected):
        summed = {ch: _summed_channel(ch, corrected) for ch in channels}
        if ratio:
            for ch in (roi_ratio_num.value, roi_ratio_den.value):
                summed.setdefault(ch, _summed_channel(ch, corrected))
        header = ["#", "ROI", "px"]
        for ch in channels:
            header += [f"{ch} tot", f"{ch} mean"]
        if ratio:
            header += [f"{roi_ratio_num.value}/{roi_ratio_den.value}"]
        rows = []
        for idx, (name, info) in enumerate(_roi["rois"].items(), 1):
            m = info["mask"]
            npix = int(m.sum())
            row = [str(idx), name, f"{npix:,}"]
            for ch in channels:
                vals = summed[ch][m] if npix else np.array([0.0])
                row += [f"{float(vals.sum()):,.0f}",
                        f"{float(vals.mean()) if npix else 0.0:,.1f}"]
            if ratio:
                num = float(summed[roi_ratio_num.value][m].sum()) if npix else 0.0
                den = float(summed[roi_ratio_den.value][m].sum()) if npix else 0.0
                row += [f"{(num / den):.4f}" if den else "—"]
            rows.append(row)
        return header, rows

    def _stats_text(header, rows, source_label):
        widths = [max(len(header[c]), *(len(r[c]) for r in rows)) if rows
                  else len(header[c]) for c in range(len(header))]
        def fmt(vals):
            return "  ".join(v.rjust(widths[c]) for c, v in enumerate(vals))
        lines = [f"ROI statistics  ·  {os.path.basename(img.path)}",
                 f"counts source: {source_label}",
                 "(totals are pooled counts over the summed stack; "
                 "ratio is pooled Σnum/Σden)", "",
                 fmt(header), "-" * (sum(widths) + 2 * (len(widths) - 1))]
        lines += [fmt(r) for r in rows]
        return "\n".join(lines)

    def _stats_fig(header, rows, source_label):
        ncol, nrow = len(header), len(rows) + 1
        w_in = min(2.4 + 1.6 * ncol, 34)
        h_in = max(1.8, 0.7 + 0.52 * nrow)
        fig, ax = plt.subplots(figsize=(w_in, h_in))
        ax.axis("off")
        tbl = ax.table(cellText=rows or [["—"] * ncol], colLabels=header,
                       loc="center", cellLoc="center")
        tbl.auto_set_font_size(False); tbl.set_fontsize(11); tbl.scale(1, 1.6)
        for (r, c), cell in tbl.get_celld().items():
            if r == 0:
                cell.set_text_props(fontweight="bold")
            if c <= 1:
                cell.set_text_props(ha="left")
        ax.set_title(f"ROI statistics ({source_label}) — "
                     f"{os.path.basename(img.path)}", fontsize=12, pad=14)
        fig.tight_layout()
        return fig

    def _roi_compute(event):
        if not _roi["rois"]:
            roi_status.object = "⚠️ No ROIs — add some from clusters or detection."
            return
        channels = list(roi_stat_chans.value) or [non_se[0]]
        ratio = bool(roi_ratio_on.value)
        corrected = roi_counts_src.value.startswith("Drift")
        source_label = ("drift-corrected" if (corrected and img.corrected is not None)
                        else "raw integer counts" if not corrected
                        else "raw (no drift applied)")
        try:
            header, rows = _build_stats_rows(channels, ratio, corrected)
        except Exception as exc:
            roi_status.object = f"❌ Stats failed: `{exc}`"
            return
        _roi["stats_text"] = _stats_text(header, rows, source_label)
        fig = _stats_fig(header, rows, source_label)
        w_in, h_in = fig.get_size_inches()
        dpi = 150
        px_w, px_h = int(w_in * dpi), int(h_in * dpi)
        _roi["stats_png"] = _fig_to_png(fig, dpi=dpi)   # closes the figure
        objs = []
        objs.append(pn.Column(_html_table(header, rows, width=920, right_from=2),
                              scroll=True, sizing_mode="stretch_width"))
        # Journal-preview render at natural (large) size in a scroll box —
        # horizontal scroll for wide many-channel tables.
        objs.append(pn.pane.Markdown(
            "<span style='font-size:11px;opacity:0.6'>Journal preview "
            "(scroll if wide):</span>"))
        objs.append(pn.Column(pn.pane.PNG(_roi["stats_png"], width=px_w),
                              scroll=True, height=min(px_h + 24, 460),
                              sizing_mode="stretch_width"))
        roi_table_box.objects = objs
        # refresh the isolate dropdown and draw the numbered map
        roi_show.options = ["all"] + list(_roi["rois"].keys())
        if roi_show.value not in roi_show.options:
            roi_show.value = "all"
        _render_map()
        roi_status.object = (f"✅ Stats for {len(_roi['rois'])} ROI(s) over "
                             f"{len(channels)} channel(s) · {source_label}.")
    roi_compute_btn.on_click(_roi_compute)

    def _roi_remove(event):
        if _roi["rois"]:
            last = list(_roi["rois"])[-1]
            _roi["rois"].pop(last)
            _refresh_roi_list()
            roi_status.object = f"Removed `{last}` — {len(_roi['rois'])} remain."
        else:
            roi_status.object = "No ROIs to remove."
    roi_remove_btn.on_click(_roi_remove)

    def _roi_clear(event):
        _roi["rois"].clear()
        _refresh_roi_list()
        roi_overlay_box.objects = [pn.pane.Markdown(
            "<span style='opacity:0.6'>No detection yet.</span>")]
        roi_table_box.objects = [pn.pane.Markdown(
            "<span style='opacity:0.6'>No stats yet.</span>")]
        roi_status.object = "Cleared all ROIs."
    roi_clear_btn.on_click(_roi_clear)

    def _roi_add_fig(event):
        if _roi["stats_png"] is None:
            roi_status.object = "Nothing to add — compute stats first."
            return
        cap = (f"ROI statistics, {len(_roi['rois'])} ROI(s)  ·  "
               f"{os.path.basename(img.path)}")
        JOURNAL.append({"png": _roi["stats_png"], "caption": cap,
                        "added": datetime.datetime.now().strftime("%Y-%m-%d %H:%M")})
        _update_journal_status()
        roi_status.object = f"✅ Table added — {len(JOURNAL)} item(s) in journal."
    roi_add_fig_btn.on_click(_roi_add_fig)

    def _roi_add_txt(event):
        if not _roi["stats_text"]:
            roi_status.object = "Nothing to add — compute stats first."
            return
        cap = f"ROI statistics summary  ·  {os.path.basename(img.path)}"
        _journal_add_text(_roi["stats_text"], cap)
        roi_status.object = f"✅ Summary added — {len(JOURNAL)} item(s) in journal."
    roi_add_txt_btn.on_click(_roi_add_txt)

    # -- Depth profiles (signal vs plane / sputter depth) --------------------
    n_planes = int(img.metadata.get("n_planes", 1))
    dp_chan = pn.widgets.Select(name="Profile channel", options=list(masses),
                                value=non_se[0], width=160)
    dp_metric = pn.widgets.Select(name="Metric",
                                  options=["mean counts / px", "total counts"],
                                  value="mean counts / px", width=170)
    dp_ratio = pn.widgets.Checkbox(name="Ratio profile (uses ratio num/den above)",
                                   value=False)
    dp_btn = pn.widgets.Button(name="Plot depth profiles", button_type="primary",
                               icon="chart-line", width=180)
    dp_add_fig_btn = pn.widgets.Button(name="➕ Add profile to journal",
                                       button_type="default", width=210)
    dp_box = pn.Column(pn.pane.Markdown(
        "<span style='opacity:0.6'>No profile yet — add ROIs and click "
        "**Plot depth profiles**.</span>"))
    _dp = {"png": None}

    def _dp_series(channel, mask, corrected):
        """Per-plane pooled counts within a mask → 1-D array (length n_planes)."""
        if corrected:
            stack = np.asarray(img.get_channel(channel))      # (planes, H, W)
        else:
            ch = list(img.masses).index(channel)
            stack = img.data[:, ch, :, :].astype(float)
        return stack[:, mask].sum(axis=1).astype(float)

    def _dp_plot():
        rois = _roi["rois"]
        if not rois:
            return None
        corrected = roi_counts_src.value.startswith("Drift")
        ratio = bool(dp_ratio.value)
        metric_mean = dp_metric.value.startswith("mean")
        cmap = plt.get_cmap("tab10")
        fig, ax = plt.subplots(figsize=(8.0, 4.8))
        x = np.arange(n_planes)
        for i, (name, info) in enumerate(rois.items()):
            m = info["mask"]
            npix = int(m.sum()) or 1
            if ratio:
                num = _dp_series(roi_ratio_num.value, m, corrected)
                den = _dp_series(roi_ratio_den.value, m, corrected)
                y = np.divide(num, den, out=np.zeros_like(num),
                              where=den > 0)
                ylabel = f"{roi_ratio_num.value} / {roi_ratio_den.value}  (pooled)"
            else:
                tot = _dp_series(dp_chan.value, m, corrected)
                y = (tot / npix) if metric_mean else tot
                ylabel = (f"{dp_chan.value}  "
                          + ("mean counts / px" if metric_mean else "total counts"))
            xx = x if len(y) == len(x) else np.arange(len(y))
            ax.plot(xx, y, marker="o", ms=3, lw=1.3, color=cmap(i % 10),
                    label=f"{name}  ({npix:,} px)")
        src = "drift-corrected" if corrected and img.corrected is not None else \
              ("raw" if not corrected else "raw (no drift)")
        ax.set_xlabel("plane index  (sputter depth →)")
        ax.set_ylabel(ylabel)
        ax.set_title(f"ROI depth profiles ({src}) — {os.path.basename(img.path)}",
                     fontsize=10)
        ax.legend(fontsize=8, ncol=2, framealpha=0.4)
        ax.grid(alpha=0.25)
        fig.tight_layout()
        return fig

    def _dp_run(event):
        if not _roi["rois"]:
            roi_status.object = "⚠️ No ROIs — add some before plotting profiles."
            return
        try:
            fig = _dp_plot()
        except Exception as exc:
            roi_status.object = f"❌ Depth profile failed: `{exc}`"
            return
        _dp["png"] = _fig_to_png(fig)
        dp_box.objects = [pn.pane.PNG(_dp["png"], width=760)]
        roi_status.object = (f"✅ Depth profiles for {len(_roi['rois'])} ROI(s) "
                             f"over {n_planes} planes.")
    dp_btn.on_click(_dp_run)

    def _dp_add_fig(event):
        if _dp["png"] is None:
            roi_status.object = "Nothing to add — plot a profile first."
            return
        what = (f"{roi_ratio_num.value}/{roi_ratio_den.value} ratio"
                if dp_ratio.value else dp_chan.value)
        cap = (f"ROI depth profiles — {what}, {len(_roi['rois'])} ROI(s)  ·  "
               f"{os.path.basename(img.path)}")
        JOURNAL.append({"png": _dp["png"], "caption": cap,
                        "added": datetime.datetime.now().strftime("%Y-%m-%d %H:%M")})
        _update_journal_status()
        roi_status.object = f"✅ Profile added — {len(JOURNAL)} item(s) in journal."
    dp_add_fig_btn.on_click(_dp_add_fig)

    roi_status = pn.pane.Markdown("")
    roi_list = pn.pane.Markdown(
        "<span style='opacity:0.6'>No ROIs yet.</span>")
    roi_note = pn.pane.Markdown(
        "<span style='font-size:11px;opacity:0.75'>ROI counts are pooled over "
        "the summed stack. <b>Counts source</b>: <i>drift-corrected</i> uses the "
        "registered float stack (matches what clustering saw); <i>raw counts</i> "
        "pulls true integer detected-ion counts from the raw stack, for Poisson "
        "error propagation. Caveat — masks are defined in drift-corrected "
        "geometry, so raw totals over a mask are approximate when drift is large "
        "(the unregistered planes smear the boundary). Isotope ratios use pooled "
        "Σnumerator / Σdenominator over each ROI — the analytically correct way, "
        "not a mean of per-pixel ratios. <b>Depth profiles</b> apply each ROI "
        "mask plane-by-plane down the stack (x = plane index; no per-plane depth "
        "calibration is in the header). The current ROI set is attached to the "
        "loaded image as <code>img.rois</code>. Edge/threshold detection works "
        "on one channel; cluster ROIs come straight from the clustering run "
        "above.</span>")

    roi_section = pn.Column(
        pn.pane.Markdown("**ROI Manager**  <span style='font-size:11px;"
                         "opacity:0.7'>(clusters → ROIs, edge detection, pooled "
                         "counts &amp; ratios, depth profiles)</span>"),
        pn.pane.Markdown("<span style='font-size:12px;opacity:0.8'>"
                         "**From clusters**</span>"),
        pn.Row(roi_load_btn, roi_cluster_pick, roi_add_clusters_btn),
        pn.pane.Markdown("<span style='font-size:12px;opacity:0.8'>"
                         "**From the Draw ROIs tab**</span>"),
        pn.Row(roi_load_drawn_btn),
        pn.pane.Markdown("<span style='font-size:12px;opacity:0.8'>"
                         "**Edge / threshold detection**</span>"),
        pn.Row(roi_edge_chan, roi_edge_method, roi_edge_sigma),
        pn.Row(roi_edge_pct, roi_edge_minarea, roi_edge_merge, roi_detect_btn),
        roi_overlay_box,
        pn.pane.Markdown("<span style='font-size:12px;opacity:0.8'>"
                         "**Statistics**</span>"),
        pn.Row(roi_stat_chans,
               pn.Column(roi_counts_src, roi_ratio_on),
               pn.Column(roi_ratio_num, roi_ratio_den)),
        pn.Row(roi_compute_btn, roi_remove_btn, roi_clear_btn),
        pn.Row(roi_add_fig_btn, roi_add_txt_btn),
        pn.pane.Markdown("<span style='font-size:12px;opacity:0.8'>"
                         "**Depth profiles**</span>"),
        pn.Row(dp_chan, dp_metric, dp_ratio),
        pn.Row(dp_btn, dp_add_fig_btn),
        dp_box,
        roi_list,
        roi_status,
        roi_table_box,
        pn.pane.Markdown("<span style='font-size:12px;opacity:0.8'>"
                         "**ROI map**  <span style='font-size:11px;opacity:0.7'>"
                         "(numbers match the table's # column; use *Show on map* "
                         "to isolate one)</span></span>"),
        pn.Row(roi_map_chan, roi_show),
        roi_map_box,
        roi_note,
    )

    children = [pn.pane.Markdown("### Analysis")]
    if drift_hint:
        children.append(pn.pane.Markdown(drift_hint))
    children += [clustering_section, pn.layout.Divider(),
                 roi_section, pn.layout.Divider(), distributions_section]
    return pn.Column(*children)


# ─────────────────────────────────────────────────────────────────────────────
# Draw ROIs tab  (interactive Bokeh canvas: freehand / polygon / rect / ellipse)
# ─────────────────────────────────────────────────────────────────────────────
#
# A separate INTERACTIVE tab (the Imaging tab is static matplotlib, hence the
# split). Shapes are drawn with Bokeh edit tools, rasterised to boolean masks at
# image resolution, then staged onto img._pending_drawn via "Send to Analysis".
# The Analysis ROI Manager pulls them in as a third ROI source ("drawn") and
# reuses its existing pooled-count / ratio / depth-profile / journal engine —
# no duplicated stats code lives here.
#
# Coordinate convention: the background array is displayed np.flipud'd so row 0
# sits at the top (imshow convention). Drawn vertices come back in that flipped
# data space, are rasterised there, then flipped back once so the masks match
# the data orientation that cluster / edge ROIs (and img.rois) already use.

def draw_roi_view(img):
    if img is None:
        return pn.pane.Markdown(
            "### Draw ROIs\nLoad a `.im` file to draw regions of interest "
            "by hand.")

    from bokeh.plotting import figure
    from bokeh.models import (ColumnDataSource, FreehandDrawTool, PolyDrawTool,
                              BoxEditTool, LinearColorMapper)
    from bokeh.palettes import Greys256
    from skimage.draw import polygon2mask, ellipse as sk_ellipse

    masses = img.masses
    non_se = [m for m in masses if "SE" not in m.upper()] or masses

    def _fig_to_png(fig, dpi=150):
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight",
                    facecolor=fig.get_facecolor())
        plt.close(fig)
        return buf.getvalue()

    def _summed(channel, corrected):
        """Summed-over-planes (H, W) counts — mirrors the Analysis ROI manager's
        _summed_channel so the canvas geometry matches what the stats see."""
        if corrected:
            stack = np.asarray(img.get_channel(channel))
            return (stack.sum(axis=0) if stack.ndim == 3 else stack).astype(float)
        ch = list(masses).index(channel)
        return img.data[:, ch, :, :].sum(axis=0).astype(float)

    # -- controls (consistent with the Analysis ROI stats toggles) -----------
    chan = pn.widgets.Select(name="Background channel", options=list(masses),
                             value=non_se[0], width=180)
    counts_src = pn.widgets.Select(
        name="Counts source",
        options=["Drift-corrected (if run)", "Raw counts (integer)"],
        value="Drift-corrected (if run)", width=210)

    def _corrected_now():
        return counts_src.value.startswith("Drift")

    # -- interactive canvas ---------------------------------------------------
    data0 = _summed(chan.value, _corrected_now())
    H, W = data0.shape
    vmax0 = float(np.percentile(data0, 99.5)) or 1.0

    # Bokeh's Greys256 runs black→white with increasing index (index 0 = black),
    # the same direction as matplotlib "gray": low counts dark, high counts bright.
    # (Do NOT reverse it — that inverts contrast relative to the Imaging tab.)
    pal = list(Greys256)
    cmapper = LinearColorMapper(palette=pal, low=0.0, high=vmax0)

    p = figure(width=560, height=560, match_aspect=True,
               x_range=(0, W), y_range=(0, H),
               tools="pan,wheel_zoom,reset,save", toolbar_location="above",
               title=f"Draw on {chan.value} — origin top-left")
    p.grid.visible = False
    p.xaxis.axis_label = "px (x)"
    p.yaxis.axis_label = "px (y)"
    img_glyph = p.image(image=[np.flipud(np.clip(data0, 0, vmax0))],
                        x=0, y=0, dw=W, dh=H, color_mapper=cmapper)

    free_src = ColumnDataSource(data=dict(xs=[], ys=[]))
    poly_src = ColumnDataSource(data=dict(xs=[], ys=[]))
    rect_src = ColumnDataSource(data=dict(x=[], y=[], width=[], height=[]))
    ell_src  = ColumnDataSource(data=dict(x=[], y=[], width=[], height=[]))

    free_r = p.multi_line(xs="xs", ys="ys", source=free_src,
                          line_color="#00e5ff", line_width=2)
    poly_r = p.patches(xs="xs", ys="ys", source=poly_src,
                       fill_alpha=0.12, fill_color="#ffd54f",
                       line_color="#ffd54f", line_width=2)
    rect_r = p.rect(x="x", y="y", width="width", height="height", source=rect_src,
                    fill_alpha=0.10, fill_color="#69f0ae",
                    line_color="#69f0ae", line_width=2)
    # Bokeh has no native ellipse-draw, so BoxEditTool edits a (dashed) Rect and
    # we render the inscribed ellipse from the same source for live feedback.
    ell_box_r = p.rect(x="x", y="y", width="width", height="height", source=ell_src,
                       fill_alpha=0.0, line_color="#ff8a80", line_dash="dashed",
                       line_width=1)
    p.ellipse(x="x", y="y", width="width", height="height", source=ell_src,
              fill_alpha=0.12, fill_color="#ff8a80", line_color="#ff8a80",
              line_width=2)

    free_tool = FreehandDrawTool(renderers=[free_r], num_objects=0)
    poly_tool = PolyDrawTool(renderers=[poly_r], num_objects=0)
    rect_tool = BoxEditTool(renderers=[rect_r], num_objects=0)
    ell_tool  = BoxEditTool(renderers=[ell_box_r], num_objects=0)
    free_tool.description = "Freehand lasso"
    poly_tool.description = "Click-polygon (double-click / Esc to finish)"
    rect_tool.description = "Rectangle (drag; Shift-drag to add more)"
    ell_tool.description  = "Ellipse — drag a box, ellipse inscribed inside it"
    p.add_tools(free_tool, poly_tool, rect_tool, ell_tool)
    p.toolbar.active_drag = free_tool

    canvas = pn.pane.Bokeh(p)

    def _refresh_bg(*_):
        data = _summed(chan.value, _corrected_now())
        vmax = float(np.percentile(data, 99.5)) or 1.0
        cmapper.high = vmax
        img_glyph.data_source.data = dict(
            image=[np.flipud(np.clip(data, 0, vmax))])
        p.title.text = f"Draw on {chan.value} — origin top-left"
    chan.param.watch(_refresh_bg, "value")
    counts_src.param.watch(_refresh_bg, "value")

    # -- rasterise → masks ----------------------------------------------------
    _draw = {"masks": {}}
    preview_box = pn.Column(pn.pane.Markdown(
        "<span style='opacity:0.6'>Nothing rasterised yet.</span>"))
    status = pn.pane.Markdown("")

    def _overlay_fig(masks):
        data = _summed(chan.value, _corrected_now())
        vmax = float(np.percentile(data, 99.5)) or 1.0
        fig, ax = plt.subplots(figsize=(5.8, 5.8))
        ax.imshow(data, cmap="gray", vmin=0, vmax=vmax, interpolation="nearest")
        cmap = plt.get_cmap("tab10")
        for i, (name, m) in enumerate(masks.items()):
            ax.contour(m.astype(float), levels=[0.5],
                       colors=[cmap(i % 10)], linewidths=1.3)
            ys_, xs_ = np.where(m)
            if len(xs_):
                ax.text(xs_.mean(), ys_.mean(), name, color=cmap(i % 10),
                        fontsize=8, ha="center", va="center")
        ax.set_title(f"{len(masks)} drawn ROI(s) on {chan.value}", fontsize=10)
        ax.set_xlabel("px"); ax.set_ylabel("px")
        fig.tight_layout()
        return fig

    def _compute_masks():
        """Read the live edit-tool sources and rasterise every shape to an
        (H, W) bool mask in data orientation. Returns (masks, n_raw) where
        n_raw is how many raw shapes were read back from the canvas — if that
        is 0 while shapes are clearly drawn, the client→server CDS sync didn't
        land (vs. n_raw > 0 with empty masks, which is a geometry problem)."""
        masks, n_raw = {}, 0

        def _add(name, m_disp):
            m = np.flipud(np.asarray(m_disp, bool))   # flip back to data orient.
            if m.any():
                masks[name] = m

        # freehand + click-polygon: per-shape vertex lists in (xs, ys)
        for src, prefix in ((free_src, "lasso"), (poly_src, "poly")):
            xs_all = list(src.data.get("xs", []) or [])
            ys_all = list(src.data.get("ys", []) or [])
            for j, (xs, ys) in enumerate(zip(xs_all, ys_all), 1):
                n_raw += 1
                if xs is None or len(xs) < 3:
                    continue
                verts = np.column_stack([np.asarray(ys, float),
                                         np.asarray(xs, float)])   # (row, col)
                _add(f"{prefix}_{j}", polygon2mask((H, W), verts))

        # rectangles (centre x/y + width/height from BoxEditTool)
        xs = list(rect_src.data.get("x", []) or []); ys = list(rect_src.data.get("y", []) or [])
        ws = list(rect_src.data.get("width", []) or []); hs = list(rect_src.data.get("height", []) or [])
        for j, (cx, cy, w, h) in enumerate(zip(xs, ys, ws, hs), 1):
            n_raw += 1
            if not w or not h:
                continue
            r0 = int(np.floor(max(0.0, cy - abs(h) / 2)))
            r1 = int(np.ceil(min(float(H), cy + abs(h) / 2)))
            c0 = int(np.floor(max(0.0, cx - abs(w) / 2)))
            c1 = int(np.ceil(min(float(W), cx + abs(w) / 2)))
            m = np.zeros((H, W), bool); m[r0:r1, c0:c1] = True
            _add(f"rect_{j}", m)

        # ellipses inscribed in the drawn box
        xs = list(ell_src.data.get("x", []) or []); ys = list(ell_src.data.get("y", []) or [])
        ws = list(ell_src.data.get("width", []) or []); hs = list(ell_src.data.get("height", []) or [])
        for j, (cx, cy, w, h) in enumerate(zip(xs, ys, ws, hs), 1):
            n_raw += 1
            if not w or not h:
                continue
            rr, cc = sk_ellipse(cy, cx, abs(h) / 2, abs(w) / 2, shape=(H, W))
            m = np.zeros((H, W), bool); m[rr, cc] = True
            _add(f"ell_{j}", m)

        return masks, n_raw

    def _preview(masks):
        png = _fig_to_png(_overlay_fig(masks))
        preview_box.objects = [pn.pane.PNG(png, width=560)]

    def _empty_msg(n_raw):
        if n_raw == 0:
            return ("Canvas reads 0 shapes — the drawing didn't sync back. "
                    "Click once on the canvas (or pan a hair) after drawing, "
                    "then try again.")
        return (f"Read {n_raw} shape(s) but none rasterised — lasso/polygon "
                f"need \u22653 points and boxes need non-zero size.")

    def _rasterise(event):
        try:
            masks, n_raw = _compute_masks()
        except Exception as exc:
            status.object = f"\u274c Rasterise failed: `{exc}`"
            return
        _draw["masks"] = masks
        if not masks:
            preview_box.objects = [pn.pane.Markdown(
                f"<span style='opacity:0.6'>{_empty_msg(n_raw)}</span>")]
            status.object = "Nothing to rasterise."
            return
        _preview(masks)
        total = sum(int(m.sum()) for m in masks.values())
        status.object = (f"\u2705 Rasterised {len(masks)} shape(s) "
                         f"(from {n_raw} drawn), {total:,} px total. "
                         f"Click **\u25b6 Send to Analysis** to stage them.")

    # -- buttons --------------------------------------------------------------
    raster_btn = pn.widgets.Button(name="▦ Rasterise (preview)",
                                   button_type="primary", width=200)
    send_btn = pn.widgets.Button(name="▶ Send to Analysis",
                                 button_type="success", width=190)
    clear_btn = pn.widgets.Button(name="Clear canvas", button_type="default",
                                  width=140)
    raster_btn.on_click(_rasterise)

    def _send(event):
        # Always recompute from the live canvas so a forgotten Rasterise click
        # can't strand the user with a stale/empty staging set.
        try:
            masks, n_raw = _compute_masks()
        except Exception as exc:
            status.object = f"\u274c Send failed during rasterise: `{exc}`"
            return
        _draw["masks"] = masks
        if not masks:
            status.object = f"\u26a0\ufe0f Nothing to send — {_empty_msg(n_raw)}"
            return
        _preview(masks)
        img._pending_drawn = {k: v.copy() for k, v in masks.items()}
        status.object = (f"\u2705 Staged {len(masks)} ROI(s). In the "
                         f"**Analysis** tab \u2192 ROI Manager, click "
                         f"**\u2b07 Load drawn ROIs from Draw tab**.")
    send_btn.on_click(_send)

    def _clear(event):
        free_src.data = dict(xs=[], ys=[])
        poly_src.data = dict(xs=[], ys=[])
        rect_src.data = dict(x=[], y=[], width=[], height=[])
        ell_src.data  = dict(x=[], y=[], width=[], height=[])
        _draw["masks"] = {}
        preview_box.objects = [pn.pane.Markdown(
            "<span style='opacity:0.6'>Canvas cleared.</span>")]
        status.object = "Canvas cleared."
    clear_btn.on_click(_clear)

    note = pn.pane.Markdown(
        "<span style='font-size:11px;opacity:0.75'>Pick a tool from the canvas "
        "toolbar: <b>freehand lasso</b>, <b>click-polygon</b> (double-click or "
        "Esc to close), <b>rectangle</b>, or <b>ellipse</b>. Bokeh has no native "
        "ellipse-draw, so the ellipse tool edits a box and inscribes the ellipse "
        "within it (dashed box = bounds). <b>Rasterise</b> snapshots everything "
        "currently on the canvas into masks at image resolution and previews "
        "them; <b>Send to Analysis</b> stages them as <code>source: drawn</code>. "
        "Then go to <b>Analysis ▸ ROI Manager ▸ Load drawn ROIs</b> to pull them "
        "into the same stats / ratio / depth-profile / journal engine the cluster "
        "and edge ROIs use. Canvas, masks and Analysis stats must share one "
        "binning — draw, send and compute without changing the global bin in "
        "between.</span>")

    return pn.Column(
        pn.pane.Markdown("### Draw ROIs"),
        pn.Row(chan, counts_src),
        canvas,
        pn.Row(raster_btn, send_btn, clear_btn),
        status,
        pn.pane.Markdown("<span style='font-size:12px;opacity:0.8'>"
                         "**Rasterised preview**</span>"),
        preview_box,
        note,
    )


# ─────────────────────────────────────────────────────────────────────────────
# 3D tab  (interactive point cloud / isosurface via plotly) + stack export
# ─────────────────────────────────────────────────────────────────────────────

def three_d_view(img):
    if img is None:
        return pn.pane.Markdown(
            "### 3D\nLoad a `.im` file to render the stack in 3D.")

    masses = img.masses
    non_se = [m for m in masses if "SE" not in m.upper()] or masses
    n_planes = int(img.metadata.get("n_planes", 1))
    field = float(img.metadata.get("field_um", 1.0)) or 1.0

    def _stack(corrected):
        if corrected:
            return np.asarray(img.get_channel(channel.value))        # (P, H, W)
        ch = list(masses).index(channel.value)
        return img.data[:, ch, :, :].astype(float)

    # ── controls ─────────────────────────────────────────────────────────────
    mode = pn.widgets.RadioButtonGroup(
        name="Mode", options=["Point cloud", "Isosurface", "Shells (points)"],
        value="Point cloud", button_type="primary")
    channel = pn.widgets.Select(name="Channel", options=list(masses),
                                value=non_se[0], width=150)
    src = pn.widgets.Select(name="Counts source",
                            options=["Drift-corrected (if run)",
                                     "Raw counts (integer)"],
                            value="Drift-corrected (if run)", width=200)
    cscale = pn.widgets.Select(
        name="Colour scale",
        options=["Viridis", "Magma", "Inferno", "Cividis", "Hot", "Greys"],
        value="Viridis", width=130)
    log_scale = pn.widgets.Checkbox(name="Log colour scale", value=False)
    z_aspect = pn.widgets.FloatSlider(name="Z aspect (visual)", start=0.1,
                                      end=5.0, step=0.1, value=1.0, width=200)
    # point-cloud-specific
    thr_pct = pn.widgets.FloatSlider(name="Intensity threshold (percentile)",
                                     start=0.0, end=99.9, step=0.5, value=85.0,
                                     width=260)
    max_pts = pn.widgets.IntInput(name="Max points", value=60000, start=2000,
                                  end=400000, width=130)
    msize = pn.widgets.FloatSlider(name="Marker size", start=1.0, end=8.0,
                                   step=0.5, value=3.0, width=180)
    # isosurface-specific
    iso_n = pn.widgets.IntInput(name="Surfaces", value=3, start=1, end=10,
                                width=100)
    iso_op = pn.widgets.FloatSlider(name="Opacity", start=0.05, end=1.0,
                                    step=0.05, value=0.5, width=180)
    iso_cap = pn.widgets.IntInput(name="Max grid pts (decimate)", value=200000,
                                  start=20000, end=1000000, width=180)

    render_btn = pn.widgets.Button(name="Render 3D", button_type="primary",
                                   icon="cube", width=140)
    html_dir = pn.widgets.TextInput(name="Save to (folder)",
                                    value=os.path.expanduser("~"), width=240)
    html_name = pn.widgets.TextInput(name="HTML name", value="pymims_3d",
                                     width=160)
    html_btn = pn.widgets.Button(name="💾 Save interactive HTML",
                                 button_type="default", width=210)
    status = pn.pane.Markdown("")
    plot_box = pn.Column(pn.pane.Markdown(
        "<span style='opacity:0.6'>Set options and click **Render 3D**. "
        "Point cloud shows voxels above the threshold; isosurface draws "
        "nested shells through the volume.</span>"))
    _fig = {"obj": None}

    def _colorscale():
        return cscale.value

    def _log_ticks(cmin, cmax):
        """Nice 1/2/3/5×10^k count values within [cmin, cmax], as (log10
        positions, count labels) so a log colourbar still reads in counts."""
        cands = [m * (10.0 ** k) for k in range(-1, 8) for m in (1, 2, 3, 5)]
        sel = [c for c in cands if cmin <= c <= cmax]
        if len(sel) < 2:
            sel = [max(cmin, 1.0), max(cmax, cmin + 1.0)]
        return [float(np.log10(c)) for c in sel], [f"{c:g}" for c in sel]

    def _color_and_bar(count_vals):
        """Return (values-for-colour, colorbar-dict). Log mode maps colour to
        log10(counts) but labels the bar in real counts."""
        cv = np.asarray(count_vals, dtype=float)
        if log_scale.value:
            c = np.log10(np.clip(cv, 1.0, None))
            lo = max(float(np.nanmin(cv[cv > 0])) if (cv > 0).any() else 1.0, 1.0)
            tv, tt = _log_ticks(lo, float(np.nanmax(cv)))
            return c, dict(title="counts (log)", tickvals=tv, ticktext=tt)
        return cv, dict(title="counts")

    def _build_pointcloud():
        import plotly.graph_objects as go
        corrected = src.value.startswith("Drift")
        s = _stack(corrected)
        P, Hh, Ww = s.shape
        thr = float(np.percentile(s, float(thr_pct.value)))
        sel = s > thr
        n_sel = int(sel.sum())
        if n_sel == 0:
            return None, "Nothing above threshold — lower the percentile."
        zz, yy, xx = np.where(sel)
        vals = s[sel]
        cap = int(max_pts.value)
        note = ""
        if n_sel > cap:                       # random subsample for responsiveness
            rng = np.random.default_rng(0)
            idx = rng.choice(n_sel, size=cap, replace=False)
            zz, yy, xx, vals = zz[idx], yy[idx], xx[idx], vals[idx]
            note = f" (showing {cap:,} of {n_sel:,} voxels)"
        x_um = xx / max(Ww - 1, 1) * field
        y_um = yy / max(Hh - 1, 1) * field
        cvals, cbar = _color_and_bar(vals)
        fig = go.Figure(go.Scatter3d(
            x=x_um.tolist(), y=y_um.tolist(), z=zz.astype(float).tolist(),
            mode="markers", customdata=vals.tolist(),
            marker=dict(size=float(msize.value), color=cvals.tolist(),
                        colorscale=_colorscale(), opacity=1.0,
                        line=dict(width=0), colorbar=cbar),
            hovertemplate="x=%{x:.2f}µm<br>y=%{y:.2f}µm<br>"
                          "plane=%{z:.0f}<br>counts=%{customdata:.0f}"
                          "<extra></extra>"))
        src_lbl = ("drift-corrected" if corrected and img.corrected is not None
                   else "raw")
        fig.update_layout(
            scene=dict(xaxis_title="µm", yaxis_title="µm",
                       zaxis_title="plane (depth →)",
                       aspectmode="manual",
                       aspectratio=dict(x=1, y=1, z=float(z_aspect.value))),
            margin=dict(l=0, r=0, t=30, b=0), height=640,
            title=f"{channel.value} point cloud ({src_lbl}) · "
                  f"{os.path.basename(img.path)}")
        return fig, (f"✅ Point cloud: {channel.value}, threshold "
                     f"P{thr_pct.value:g} (>{thr:.0f} counts){note}"
                     f"{' · log colour' if log_scale.value else ''}.")

    def _build_isosurface():
        import plotly.graph_objects as go
        corrected = src.value.startswith("Drift")
        s = _stack(corrected)
        P, Hh, Ww = s.shape
        # Decimate in-plane only (keep all planes for depth fidelity).
        step = 1
        while (P * Hh * Ww) / (step * step) > int(iso_cap.value):
            step += 1
        s_d = s[:, ::step, ::step]
        Pd, Hd, Wd = s_d.shape
        zc, yc, xc = np.mgrid[0:Pd, 0:Hd, 0:Wd]
        x_um = xc / max(Wd - 1, 1) * field
        y_um = yc / max(Hd - 1, 1) * field
        z_co = zc.astype(float)   # plane index; in-plane decimation doesn't touch z
        raw = s_d.astype(float)
        if log_scale.value:
            valf = np.log10(np.clip(raw, 1.0, None))
            lo = float(raw[raw > 0].min()) if (raw > 0).any() else 1.0
            tv, tt = _log_ticks(max(lo, 1.0), float(raw.max()))
            cbar = dict(title="counts (log)", tickvals=tv, ticktext=tt)
        else:
            valf = raw
            cbar = dict(title="counts")
        vmax = float(np.percentile(valf, 99.5)) or (float(valf.max()) or 1.0)
        pos = valf[raw > 0]
        vmin = float(np.percentile(pos, 50)) if pos.size else float(valf.min())
        if vmin >= vmax:
            vmin = float(valf.min())
        fig = go.Figure(go.Isosurface(
            x=x_um.flatten().tolist(), y=y_um.flatten().tolist(),
            z=z_co.flatten().tolist(),
            value=valf.flatten().tolist(), isomin=vmin, isomax=vmax,
            surface_count=int(iso_n.value), opacity=float(iso_op.value),
            colorscale=_colorscale(),
            caps=dict(x_show=False, y_show=False, z_show=False),
            colorbar=cbar))
        src_lbl = ("drift-corrected" if corrected and img.corrected is not None
                   else "raw")
        fig.update_layout(
            scene=dict(xaxis_title="µm", yaxis_title="µm",
                       zaxis_title="plane (depth →)",
                       aspectmode="manual",
                       aspectratio=dict(x=1, y=1, z=float(z_aspect.value))),
            margin=dict(l=0, r=0, t=30, b=0), height=640,
            title=f"{channel.value} isosurface ({src_lbl}) · "
                  f"{os.path.basename(img.path)}")
        decim = f" (in-plane decimation ×{step})" if step > 1 else ""
        return fig, (f"✅ Isosurface: {channel.value}, {iso_n.value} shells"
                     f"{decim}{' · log colour' if log_scale.value else ''}.")

    def _build_shells():
        """Nested iso-shells rendered as point bands via Scatter3d — the
        rendering path that works on GPUs where Volume/Isosurface meshes don't.
        Selects voxels lying within a thin band around each of N iso-levels."""
        import plotly.graph_objects as go
        corrected = src.value.startswith("Drift")
        s = _stack(corrected)
        P, Hh, Ww = s.shape
        raw = s.astype(float)
        field_v = np.log10(np.clip(raw, 1.0, None)) if log_scale.value else raw
        pos = field_v[raw > 0]
        vmin = float(np.percentile(pos, 50)) if pos.size else float(field_v.min())
        vmax = float(np.percentile(field_v, 99.5)) or (float(field_v.max()) or 1.0)
        if vmin >= vmax:
            vmin = float(field_v.min())
        n = max(int(iso_n.value), 1)
        levels = np.linspace(vmin, vmax, n)
        spacing = (vmax - vmin) / max(n - 1, 1) if n > 1 else (vmax - vmin or 1.0)
        band = 0.12 * spacing if spacing > 0 else 0.05 * (abs(vmax) or 1.0)
        sel = np.zeros(field_v.shape, dtype=bool)
        for L in levels:
            sel |= np.abs(field_v - L) <= band
        sel &= (raw > 0)
        n_sel = int(sel.sum())
        if n_sel == 0:
            return None, "No voxels near the shell levels — try fewer surfaces."
        zz, yy, xx = np.where(sel)
        cnts = raw[sel]
        cap = int(max_pts.value)
        note = ""
        if n_sel > cap:
            rng = np.random.default_rng(0)
            idx = rng.choice(n_sel, size=cap, replace=False)
            zz, yy, xx, cnts = zz[idx], yy[idx], xx[idx], cnts[idx]
            note = f" (showing {cap:,} of {n_sel:,})"
        x_um = xx / max(Ww - 1, 1) * field
        y_um = yy / max(Hh - 1, 1) * field
        cvals, cbar = _color_and_bar(cnts)
        fig = go.Figure(go.Scatter3d(
            x=x_um.tolist(), y=y_um.tolist(), z=zz.astype(float).tolist(),
            mode="markers", customdata=cnts.tolist(),
            marker=dict(size=float(msize.value), color=cvals.tolist(),
                        colorscale=_colorscale(), opacity=1.0,
                        line=dict(width=0), colorbar=cbar),
            hovertemplate="x=%{x:.2f}µm<br>y=%{y:.2f}µm<br>"
                          "plane=%{z:.0f}<br>counts=%{customdata:.0f}"
                          "<extra></extra>"))
        src_lbl = ("drift-corrected" if corrected and img.corrected is not None
                   else "raw")
        fig.update_layout(
            scene=dict(xaxis_title="µm", yaxis_title="µm",
                       zaxis_title="plane (depth →)", aspectmode="manual",
                       aspectratio=dict(x=1, y=1, z=float(z_aspect.value))),
            margin=dict(l=0, r=0, t=30, b=0), height=640,
            title=f"{channel.value} shells ({src_lbl}) · "
                  f"{os.path.basename(img.path)}")
        return fig, (f"✅ Shells: {channel.value}, {n} level(s){note}"
                     f"{' · log colour' if log_scale.value else ''}.")

    def _render(event):
        try:
            import plotly.graph_objects as go  # noqa: F401
        except Exception:
            plot_box.objects = [pn.pane.Alert(
                "plotly is not installed. Install with "
                "`pip install --break-system-packages plotly`.",
                alert_type="warning")]
            return
        status.object = "Rendering…"
        try:
            if mode.value == "Point cloud":
                fig, msg = _build_pointcloud()
            elif mode.value == "Isosurface":
                fig, msg = _build_isosurface()
            else:
                fig, msg = _build_shells()
        except Exception as exc:
            plot_box.objects = [pn.pane.Alert(f"Render failed: `{exc}`",
                                              alert_type="danger")]
            status.object = ""
            return
        if fig is None:
            status.object = f"⚠️ {msg}"
            return
        _fig["obj"] = fig
        plot_box.objects = [pn.pane.Plotly(
            fig, height=660, sizing_mode="stretch_width",
            config={"responsive": True})]
        status.object = msg
    render_btn.on_click(_render)

    def _save_html(event):
        if _fig["obj"] is None:
            status.object = "Nothing to save — render a view first."
            return
        try:
            folder = os.path.expanduser(html_dir.value or "~")
            os.makedirs(folder, exist_ok=True)
            name = (html_name.value or "pymims_3d").strip()
            if not name.lower().endswith(".html"):
                name += ".html"
            path = os.path.join(folder, name)
            _fig["obj"].write_html(path, include_plotlyjs="cdn")
            status.object = f"✅ Saved interactive view `{name}` to `{folder}`."
        except Exception as exc:
            status.object = f"❌ HTML save failed: `{exc}`"
    html_btn.on_click(_save_html)

    # ── Stack export (TIFF / NPZ) for external tools ─────────────────────────
    exp_what = pn.widgets.Select(name="Export data",
                                 options=["Drift-corrected (if run)",
                                          "Raw stack"],
                                 value="Drift-corrected (if run)", width=200)
    exp_fmt = pn.widgets.Select(name="Format",
                                options=["TIFF (ImageJ hyperstack)", "NPZ",
                                         "Both"],
                                value="TIFF (ImageJ hyperstack)", width=210)
    exp_dir = pn.widgets.TextInput(name="Export folder",
                                   value=os.path.expanduser("~"), width=240)
    exp_name = pn.widgets.TextInput(name="Export name", value="stack", width=160)
    exp_btn = pn.widgets.Button(name="📦 Export stack", button_type="primary",
                                width=160)
    exp_status = pn.pane.Markdown("")
    exp_note = pn.pane.Markdown(
        "<span style='font-size:11px;opacity:0.75'>Writes the full 4-D stack "
        "(planes × channels × Y × X). TIFF is an ImageJ-style hyperstack "
        "(axes ZCYX, channel names embedded) that opens in any volumetric "
        "viewer; NPZ keeps the array plus masses and metadata for Python. "
        "Drift-corrected export uses the registered float stack; raw export "
        "keeps the original integer counts.</span>")

    def _export_stack(event):
        corrected = exp_what.value.startswith("Drift")
        if corrected and img.corrected is not None:
            arr = np.asarray(img.corrected)
            src_lbl = "drift-corrected"
        else:
            arr = np.asarray(img.data)
            src_lbl = "raw"
        # arr shape (planes, channels, H, W) == axes ZCYX
        folder = os.path.expanduser(exp_dir.value or "~")
        base = (exp_name.value or "stack").strip()
        want_tiff = exp_fmt.value in ("TIFF (ImageJ hyperstack)", "Both")
        want_npz = exp_fmt.value in ("NPZ", "Both")
        written = []
        try:
            os.makedirs(folder, exist_ok=True)
            if want_tiff:
                try:
                    import tifffile
                except Exception:
                    exp_status.object = ("❌ TIFF needs `tifffile` — "
                                         "`pip install --break-system-packages "
                                         "tifffile`.")
                    return
                tif_path = os.path.join(folder, base + ".tif")
                tifffile.imwrite(
                    tif_path, arr, imagej=True, metadata={
                        "axes": "ZCYX",
                        "Labels": list(masses),
                        "unit": "um",
                        "spacing": 1.0,
                    })
                written.append(os.path.basename(tif_path))
            if want_npz:
                npz_path = os.path.join(folder, base + ".npz")
                np.savez_compressed(
                    npz_path, stack=arr, masses=np.array(list(masses)),
                    axes="ZCYX", field_um=field, n_planes=int(arr.shape[0]),
                    corrected=bool(corrected and img.corrected is not None))
                written.append(os.path.basename(npz_path))
        except Exception as exc:
            exp_status.object = f"❌ Export failed: `{exc}`"
            return
        shp = "×".join(str(d) for d in arr.shape)
        exp_status.object = (f"✅ Exported {', '.join(written)} "
                             f"({src_lbl}, {shp} ZCYX) to `{folder}`.")
    exp_btn.on_click(_export_stack)

    # ── layout (mode-dependent controls) ─────────────────────────────────────
    def _controls(m):
        common = pn.Row(channel, src, cscale, z_aspect, log_scale)
        if m == "Point cloud":
            return pn.Column(common, pn.Row(thr_pct, max_pts, msize))
        if m == "Isosurface":
            return pn.Column(common, pn.Row(iso_n, iso_op, iso_cap))
        return pn.Column(common, pn.Row(iso_n, max_pts, msize))  # Shells (points)

    note = pn.pane.Markdown(
        "<span style='font-size:11px;opacity:0.75'>Interactive in-browser 3D: "
        "drag to rotate, scroll to zoom. The Z axis is plane index, not "
        "calibrated depth — and remember sputter rate varies across the field, "
        "so treat this as a data-space view, not a reconstructed geometry. "
        "<b>Point cloud</b> and <b>Shells (points)</b> use the point-sprite "
        "renderer (robust everywhere); <b>Isosurface</b> uses the mesh/volume "
        "renderer, which some integrated GPUs won't draw — if it comes up empty, "
        "use Shells, which gives the same nested-shell view from points. Save an "
        "interactive HTML to share, or export the stack below for a dedicated "
        "volumetric viewer.</span>")

    export_section = pn.Column(
        pn.layout.Divider(),
        pn.pane.Markdown("**Export stack (TIFF / NPZ)**  <span style='font-size:"
                         "11px;opacity:0.7'>for external volumetric tools</span>"),
        pn.Row(exp_what, exp_fmt),
        pn.Row(exp_dir, exp_name, exp_btn),
        exp_status,
        exp_note,
    )

    return pn.Column(
        pn.pane.Markdown("### 3D"),
        mode,
        pn.panel(pn.bind(_controls, mode.param.value)),
        pn.Row(render_btn, html_dir, html_name, html_btn),
        status,
        note,
        plot_box,
        export_section,
    )


# ─────────────────────────────────────────────────────────────────────────────
# Help tab  (in-app user guide; click-to-expand sections)
# ─────────────────────────────────────────────────────────────────────────────
#
# Static content — no image needed, so it is built once and not bound to state.
# Sections are an Accordion (click to expand) rather than hover tooltips: hover
# fails on touch devices and section-level hover needs per-box CSS. Per-widget
# `description=` tooltips are the intended complementary polish layer (V1.x).

def help_view():
    def card(title, md):
        return (title, pn.pane.Markdown(md, sizing_mode="stretch_width"))

    intro = pn.pane.Markdown(
        "### Help & user guide\n"
        "PyMIMS is a toolkit for Cameca NanoSIMS `.im` files: load a stack, "
        "clean and drift-correct it, build ratio / HSI images, do high-mass-"
        "resolution deconvolution, cluster and define ROIs, profile in depth, "
        "and collect everything into an exportable journal. Expand a section "
        "below for details. A typical end-to-end run:\n\n"
        "1. **Load** a `.im` file from the sidebar.\n"
        "2. **Imaging → Plane QC**: scan planes and drop bad ones *before* "
        "drift correction.\n"
        "3. **Imaging**: run drift correction; inspect single-channel, ratio "
        "and HSI views.\n"
        "4. **Analysis**: cluster the stack and/or define ROIs (from clusters, "
        "the **Draw ROIs** tab, or edge detection).\n"
        "5. **Analysis**: compute pooled ROI stats / ratios and depth profiles.\n"
        "6. **Journal** (sidebar): collect figures and tables, then export to "
        "Word or PowerPoint.\n\n"
        "Throughout: **display filters and ratio-median are cosmetic** and never "
        "touch the data; **binning is non-destructive**; **drift-corrected vs raw "
        "counts** is a real analytical choice (see *Tips* at the bottom).",
        sizing_mode="stretch_width")

    sidebar_md = (
        "**Load a file** — set the directory, optionally drill into a subfolder, "
        "pick a `.im` file and **Load**. The summary line shows dimensions, plane "
        "count, channels and field of view.\n\n"
        "**Binning (sum-pool)** — sum-pools counts into n×n blocks and applies "
        "uniformly to *every* image (single channel, ratio, HSI) so they all "
        "co-register. Non-destructive: the full-resolution data is never altered. "
        "Use this for genuine noise reduction (it enlarges the analytical volume), "
        "unlike the display-only ratio-median.\n\n"
        "**Session** — **Save session** writes the cleaned / drift-corrected "
        "stack, the dropped-plane record, metadata and raw header to a `.npz` you "
        "can reload later without the original `.im`. **Open session** lists `.npz` "
        "files in the session folder. Note: the binning level is a *view* setting "
        "and is **not** saved with the session.\n\n"
        "**Journal** — the running collection of figures and tables you add from "
        "any tab (each **➕ Add…** button). Export the whole thing to a Word "
        "(`.docx`) or PowerPoint (`.pptx`) report. Remove-last / clear manage the "
        "contents.")

    imaging_md = (
        "Three view modes plus plane cleaning and drift correction.\n\n"
        "**Display filters** are *display only* — they don't touch the underlying "
        "data; ratios, HSI and the Analysis tab always use the raw counts.\n\n"
        "**Single channel** — pick a channel, sum all planes or scrub a single "
        "plane, choose a colormap and an upper-percentile clip.\n\n"
        "**Ratio** — numerator / denominator with optional masking: *Min counts* "
        "(denominator floor), *Max σ/R* (relative-error ceiling, Poisson), and a "
        "*δ reference* (0 = image median) for delta-notation maps.\n\n"
        "**HSI (hue-saturation-intensity)** — ratio drives hue, a chosen channel "
        "(denominator / numerator / sum) drives intensity. Includes the *Classic "
        "OpenMIMS LUT*, a scale factor, manual hue min/max, and a *ratio-median* "
        "that smooths the ratio **after** it is formed (cosmetic — for real noise "
        "reduction use the global bin).\n\n"
        "**Plane QC — run BEFORE drift correction.** Scans per-plane total counts "
        "and flags outliers (charging events, glitches, blank frames); bad planes "
        "are the main cause of silent drift-correction failures. Auto-drop "
        "(preview then apply) or drop specific indices manually. Dropping is "
        "destructive *in memory*; **Reset to raw** reloads from disk (also clears "
        "drift). To eyeball a suspect plane: Single channel → untick *Sum all "
        "planes* → scrub to it → *add scrubber plane*.\n\n"
        "**Drift correction** registers the planes; clustering, GMM fits and the "
        "drift-corrected ROI counts all use the registered stack, so run it once "
        "QC is done. **Save at DPI** exports the current figure at a chosen "
        "resolution.")

    hmr_md = (
        "High-mass-resolution (HMR) work is about telling apart ions that sit at "
        "almost exactly the same mass — for instance ¹²C¹⁵N and ¹³C¹⁴N, which an "
        "ordinary mass setting lumps together as one. This tab analyses an "
        "exported HMR scan (`.csv` / `.xlsx`).\n\n"
        "**How an HMR peak is shaped — the idea everything else rests on.** A peak "
        "here is built by scanning the secondary ion beam across the exit slit. As "
        "the beam slides *onto* the slit, the signal ramps up; while the beam sits "
        "*fully inside* the slit, the signal is flat; as it slides *off*, the "
        "signal ramps back down. So a single species doesn't make a sharp spike — "
        "it makes a **flat-topped peak with sloped sides**. Two things set that "
        "shape: the **slit width** sets how wide the flat top is, and the **beam "
        "width** (how tightly the beam is focused) sets how steep the sloped sides "
        "are. A sharp, well-focused beam gives steep edges; a fat or aberrated "
        "beam gives lazy, smeared ones. Internally each peak is modelled as a "
        "rectangle (the slit) softened by a bell curve (the beam), which is just "
        "this picture written as maths.\n\n"
        "**Spectrum** — load and view the scan; toggle the x-axis units and a "
        "log y-scale.\n\n"
        "**IRF & resolving power** — the IRF (instrument response function) is just "
        "the shape a single, isolated peak makes; everything downstream is built "
        "on it. Pick a rising or falling edge of a peak and a reference mass, and "
        "the tool measures **how steep that edge is**, reporting two numbers: the "
        "familiar **MRP** (mass resolving power, taken from the width of the "
        "10–90 % portion of the edge — bigger means you can separate closer "
        "neighbours) and the **edge σ**, a direct read on how much the beam is "
        "blurring things. Steeper edge → higher MRP → cleaner separation. (MRP "
        "only describes that 10–90 % stretch of the edge, so it says nothing about "
        "the very top, the tails, or whether the two sides are symmetric.)\n\n"
        "**Deconvolution** — when two species are so close in mass that their "
        "flat-topped peaks overlap into one lumpy hump, this step un-blurs the "
        "hump back into separate peaks. It already knows the characteristic "
        "flat-top-with-sloped-edges shape (from the edge above), then asks: *what "
        "is the smallest set of peaks of that shape, at what positions and heights, "
        "that reproduces the measured hump?* The 'smallest set' part matters — it "
        "prefers the simplest explanation rather than inventing extra peaks to "
        "chase noise. It tries a range of peak counts and keeps the best-supported "
        "one, then shows the recovered components and the **residual** (measured "
        "minus model) so you can judge the fit. The controls bound the peak-count "
        "range and set a minimum peak height, relative to the noise, below which a "
        "component is thrown out.\n\n"
        "**Species identification** — click the elements you expect in the sample "
        "on the periodic table, set the polarity and the maximum atoms per "
        "molecular ion, and the tool assigns the recovered peaks to candidate ions "
        "by their **spacing** along the mass axis. The axis is the instrument's "
        "nominal scale (not absolutely calibrated), so only the *relative* gaps "
        "between peaks are used — not their absolute positions. *Reference peak* "
        "mode: you tag one peak you're sure of and everything else is read off by "
        "mass difference from it. *Blind* mode: it guesses purely from the "
        "separation pattern, so trust it only when the next-best assignment is "
        "clearly worse. Ticking *labeled sample* stops natural isotope abundances "
        "from skewing the ranking.\n\n"
        "**Journal & export** — add the spectrum, residuals, peak table or "
        "identification to the journal.\n\n"
        "*If the tab reports the engine isn't importable, the modules `irf.py`, "
        "`deconvolve.py`, `sparse_deconvolve.py` and `hmr_identify.py` must sit in "
        "the repo root next to the app.*")

    analysis_md = (
        "**Clustering** — k-means or hierarchical over a chosen feature space "
        "(log z-scored, robust-z, log, raw) and channel set (SE excluded by "
        "default). *Min counts* and a *mask channel* exclude low-signal pixels; "
        "*subsample* speeds up hierarchical linkage. Runs over the drift-corrected "
        "stack — run drift correction first.\n\n"
        "**ROI Manager** — one ROI set, three sources:\n"
        "- **From clusters** — load a clustering run's per-cluster masks and add "
        "the ones you want.\n"
        "- **From the Draw ROIs tab** — pull in shapes you drew and staged (see "
        "the Draw ROIs section).\n"
        "- **Edge / threshold detection** — Otsu, manual percentile, or Canny+fill "
        "on one channel, with a min-area filter and optional merge.\n\n"
        "**Statistics** — per-ROI pooled counts and means over the chosen "
        "channels, plus an optional pooled isotope **ratio** (Σnumerator / "
        "Σdenominator over each ROI — the analytically correct way, not a mean of "
        "per-pixel ratios). **Counts source**: *drift-corrected* uses the "
        "registered float stack (matches what clustering saw); *raw counts* pulls "
        "true integer detected-ion counts for Poisson error propagation. Results "
        "carry a leading **#** that ties each row to the numbered **ROI map** "
        "below; *Show on map* isolates one ROI.\n\n"
        "**Depth profiles** — apply each ROI mask plane-by-plane down the stack "
        "(x = plane index; there is no per-plane depth calibration in the header). "
        "Mean-per-px or total counts, or a ratio profile.\n\n"
        "*Caveat:* masks are defined in drift-corrected geometry, so raw totals "
        "over a mask are approximate when drift is large (unregistered planes "
        "smear the boundary). The current ROI set is attached to the loaded image "
        "as `img.rois`.")

    draw_md = (
        "An interactive Bokeh canvas for hand-drawn ROIs (the Imaging tab is "
        "static, which is why drawing lives here).\n\n"
        "Pick a tool from the canvas toolbar: **freehand lasso**, **click-"
        "polygon** (double-click or Esc to close), **rectangle**, or **ellipse**. "
        "Bokeh has no native ellipse-draw, so the ellipse tool edits a box and "
        "inscribes the ellipse within it (the dashed box is the bounds).\n\n"
        "**Send to Analysis** rasterises whatever is on the canvas into masks at "
        "image resolution and stages them (it also previews them, and recomputes "
        "from the live canvas, so you can't be stranded by skipping the preview). "
        "**Rasterise (preview)** just shows the masks without staging. The status "
        "line reports how many raw shapes it read back from the canvas.\n\n"
        "Then go to **Analysis → ROI Manager → Load drawn ROIs** to bring them "
        "into the same stats / ratio / depth-profile / journal engine as the "
        "cluster and edge ROIs (they carry the source tag `drawn`).\n\n"
        "**Keep one binning** across draw → send → compute — the canvas, masks "
        "and Analysis stats must share the same global bin.")

    threed_md = (
        "Renders the stack as an interactive 3-D view (point cloud / isosurface "
        "via Plotly) for one channel. The **Export stack** controls write the "
        "stack to TIFF or NPZ for external volumetric tools.")

    meta_md = (
        "The full reverse-engineered Cameca header for the loaded file — "
        "acquisition parameters, detector setup, raster, dwell, and the analytical "
        "metadata PyMIMS parses out of the `.im` binary.")

    tips_md = (
        "**Run Plane QC before drift correction.** Bad planes (charging, glitches, "
        "blanks) are the main cause of silent drift-correction failures.\n\n"
        "**Drift-corrected vs raw counts.** For display and for matching what "
        "clustering saw, use drift-corrected. For Poisson error propagation on ROI "
        "stats, use raw counts — those are the true integer detected-ion counts.\n\n"
        "**Binning is a view setting.** It is non-destructive and is *not* saved "
        "in a session. Keep it consistent when drawing ROIs and computing their "
        "stats, since masks and counts must share one geometry.\n\n"
        "**Filters and ratio-median are cosmetic.** They change the picture, not "
        "the numbers. Any analysis reads the raw counts.\n\n"
        "**ROI numbering.** The map's numbered discs match the stats table's # "
        "column; the centroid label is most meaningful for compact ROIs (drawn / "
        "edge), less so for space-filling cluster partitions.")

    sections = pn.Accordion(
        card("Sidebar — load · binning · sessions · journal", sidebar_md),
        card("Imaging tab", imaging_md),
        card("HMR tab — high mass resolution", hmr_md),
        card("Analysis tab — clustering · ROIs · depth profiles", analysis_md),
        card("Draw ROIs tab", draw_md),
        card("3D tab", threed_md),
        card("Metadata tab", meta_md),
        card("Tips & common pitfalls", tips_md),
        active=[0], toggle=False, sizing_mode="stretch_width",
    )

    footer = pn.pane.Markdown(
        f"<span style='font-size:11px;opacity:0.6'>PyMIMS v{PYMIMS_VERSION} — "
        "NanoSIMS .im toolkit. This guide mirrors the inline notes on each tab; "
        "if behaviour and guide ever disagree, the inline note is authoritative. "
        "PyMIMS is free and open-source; if it's useful in your work you can "
        "support continued development at "
        "<a href='https://n1lab.substack.com' target='_blank'>The N=1 Lab</a>."
        "</span>", sizing_mode="stretch_width")

    return pn.Column(intro, pn.layout.Divider(), sections, footer,
                     sizing_mode="stretch_width")


# ─────────────────────────────────────────────────────────────────────────────
# Ion-optics tab  (single-scan line-spread function from an HMR peak edge)
# ─────────────────────────────────────────────────────────────────────────────
#
# Physics: one edge of an HMR peak is the beam's edge-spread function (ESF) — how
# transmission ramps as the secondary beam scans onto the exit slit. Its
# derivative is the line-spread function (LSF), the resolution kernel in the
# dispersion direction. A single edge of a rect(slit, width w) ⊛ Gaussian(beam, σ)
# peak is set by σ alone (the slit width only sets how far apart the two edges
# sit), so the LSF isolates the optics from the slit — the shape information MRP
# (a 10–90 % width) structurally cannot show: asymmetry, tails, kurtosis.
#
# Reads whatever scan is loaded in the HMR tab via HMR_SHARED (populated by the
# HMR loader), so it reuses that parser rather than re-reading the file.

def _col_to_idx(letters):
    """Spreadsheet column letters (A, B, …, AA) → 0-based index."""
    n = 0
    for ch in letters:
        n = n * 26 + (ord(ch) - 64)
    return n - 1


def _read_xlsx_minimal(path):
    """Dependency-free .xlsx reader (stdlib zipfile + ElementTree), used as a
    fallback when openpyxl isn't installed. Reads the first worksheet into a
    DataFrame with the first row as the header. Handles shared strings, inline
    strings and numbers — enough for tabular exports like HMR scans. Not a
    general xlsx reader (no formulas/date formatting/multi-sheet logic)."""
    import re as _re
    import zipfile
    import pandas as pd
    from xml.etree import ElementTree as ET

    def _ln(tag):                       # local tag name without namespace
        return tag.rsplit("}", 1)[-1]

    with zipfile.ZipFile(path) as z:
        names = z.namelist()
        shared = []
        if "xl/sharedStrings.xml" in names:
            sroot = ET.fromstring(z.read("xl/sharedStrings.xml"))
            for si in sroot:
                shared.append("".join(t.text or "" for t in si.iter()
                                      if _ln(t.tag) == "t"))
        sheets = sorted(n for n in names
                        if n.startswith("xl/worksheets/sheet") and n.endswith(".xml"))
        if not sheets:
            raise ValueError("no worksheet found in xlsx")
        wroot = ET.fromstring(z.read(sheets[0]))
        rows_out = []
        for el in wroot.iter():
            if _ln(el.tag) != "row":
                continue
            cells, maxc = {}, -1
            for c in el:
                if _ln(c.tag) != "c":
                    continue
                ref = c.get("r", "")
                mo = _re.match(r"[A-Z]+", ref)
                ci = _col_to_idx(mo.group()) if mo else (max(cells) + 1 if cells else 0)
                t = c.get("t")
                vtext = None
                for ch in c:
                    if _ln(ch.tag) == "v":
                        vtext = ch.text
                    elif _ln(ch.tag) == "is":
                        vtext = "".join(tt.text or "" for tt in ch.iter()
                                        if _ln(tt.tag) == "t")
                        t = "inlineStr"
                if vtext is None:
                    val = None
                elif t == "s":
                    val = (shared[int(vtext)] if vtext.isdigit()
                           and int(vtext) < len(shared) else vtext)
                elif t in ("str", "inlineStr"):
                    val = vtext
                else:
                    try:
                        val = float(vtext)
                    except (TypeError, ValueError):
                        val = vtext
                cells[ci] = val
                maxc = max(maxc, ci)
            rows_out.append([cells.get(i) for i in range(maxc + 1)])

    if not rows_out:
        raise ValueError("empty worksheet")
    header = [str(h) if h is not None else f"col{i}"
              for i, h in enumerate(rows_out[0])]
    width = len(header)
    data = [(r + [None] * (width - len(r)))[:width] for r in rows_out[1:]]
    df = pd.DataFrame(data, columns=header)
    for col in df.columns:                       # keep numeric columns numeric
        conv = pd.to_numeric(df[col], errors="coerce")
        if conv.notna().sum() >= max(1, len(conv) // 2):
            df[col] = conv
    return df


HMR_SHARED = {"m": None, "counts": None, "axes": None,
              "axis_label": None, "name": None}


def lsf_view():
    from scipy.ndimage import gaussian_filter1d

    intro = pn.pane.Markdown(
        "### Ion optics — line-spread function (LSF)\n"
        "One edge of an HMR peak is the beam's **edge-spread function**; its "
        "derivative is the **line-spread function**, the true resolution kernel in "
        "the dispersion direction. Because a single edge of a slit-image peak is "
        "set by the **beam blur σ** alone (the slit width only sets how far apart "
        "the two edges sit), the LSF isolates the *optics* from the *slit* — the "
        "shape, tails and asymmetry that MRP (a 10–90 % number) can't show. Load "
        "the scan you have open in the **HMR** tab, choose an edge, and compute.")

    status = pn.pane.Markdown("*No scan loaded — click below to pull the HMR scan.*")
    src_btn = pn.widgets.Button(name="↻ Load scan from HMR tab",
                                button_type="primary", width=240)
    axis_sel = pn.widgets.Select(name="Axis", options=["amu"], value="amu", width=120)
    edge_sel = pn.widgets.Select(name="Edge", options=["rising", "falling", "both"],
                                 value="both", width=120)
    smooth = pn.widgets.IntInput(name="Smoothing (pts)", value=2, start=0, end=25,
                                 width=130)
    win_lo = pn.widgets.FloatInput(name="Window min (axis units, 0=off)",
                                   value=0.0, width=210)
    win_hi = pn.widgets.FloatInput(name="Window max (0=off)", value=0.0, width=150)
    dm_as = pn.widgets.FloatInput(name="Abundance-sens. Δm (mmu)", value=5.0,
                                  start=0.0, width=200)
    compute_btn = pn.widgets.Button(name="Compute LSF", button_type="primary",
                                    width=160)
    peaks_warn = pn.pane.Markdown("")
    metrics = pn.pane.Markdown("")
    fig_box = pn.Column()

    def _fig_png(fig, dpi=140):
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight",
                    facecolor=fig.get_facecolor())
        plt.close(fig)
        return buf.getvalue()

    def _pull(event=None):
        if HMR_SHARED.get("counts") is None:
            status.object = ("⚠️ No scan found. In the **HMR** tab: ↻ Scan folder "
                             "→ ⤓ Load scan, then come back and click this.")
            return
        axes = HMR_SHARED.get("axes") or {}
        axis_sel.options = list(axes.keys()) or ["amu"]
        axis_sel.value = "amu" if "amu" in axis_sel.options else axis_sel.options[0]
        status.object = (f"✅ Using **{HMR_SHARED.get('name')}** from the HMR tab. "
                         f"Choose an edge and click **Compute LSF**.")
    src_btn.on_click(_pull)

    def _interp_cross(me, esf, level):
        """m where esf first crosses `level` (monotone-ish edge)."""
        d = esf - level
        sign = np.sign(d)
        idx = np.where(np.diff(sign) != 0)[0]
        if len(idx) == 0:
            return float(me[np.argmin(np.abs(d))])
        i = idx[0]
        x0, x1, y0, y1 = me[i], me[i + 1], esf[i], esf[i + 1]
        return float(x0 + (level - y0) * (x1 - x0) / (y1 - y0)) if y1 != y0 else float(x0)

    def _moments(x, w):
        w = np.clip(np.asarray(w, float), 0, None)
        W = w.sum()
        if W <= 0 or len(x) < 3:
            return dict(c=np.nan, fwhm=np.nan, skew=np.nan, kurt=np.nan)
        c = float((x * w).sum() / W)
        var = float(((x - c) ** 2 * w).sum() / W)
        sd = np.sqrt(var) if var > 0 else np.nan
        if sd and np.isfinite(sd) and sd > 0:
            skew = float((((x - c) / sd) ** 3 * w).sum() / W)
            kurt = float((((x - c) / sd) ** 4 * w).sum() / W - 3.0)
        else:
            skew = kurt = np.nan
        half = w.max() / 2.0
        above = np.where(w >= half)[0]
        fwhm = float(x[above[-1]] - x[above[0]]) if len(above) >= 2 else np.nan
        return dict(c=c, fwhm=abs(fwhm), skew=skew, kurt=kurt)

    def _edge(m, y, side, base, amp, ic):
        if side == "rising":
            sl = slice(0, ic + 1)
        else:
            sl = slice(ic, len(m))
        me = m[sl]
        esf = np.clip((y[sl] - base) / amp, 0.0, None)
        lsf = np.abs(np.gradient(esf, me))
        return me, esf, lsf

    def _compute(event=None):
        if HMR_SHARED.get("counts") is None:
            status.object = "⚠️ Load a scan from the HMR tab first."
            return
        try:
            axes = HMR_SHARED.get("axes") or {}
            ax = axis_sel.value if axis_sel.value in axes else next(iter(axes))
            m = np.asarray(axes[ax], float)
            y = np.asarray(HMR_SHARED["counts"], float)
            order = np.argsort(m)
            m, y = m[order], y[order]
            # count peaks on the FULL scan first (drives the multi-peak warning)
            ys_full = gaussian_filter1d(y, float(smooth.value)) if smooth.value else y.copy()
            bf = float(np.percentile(ys_full, 5)); af = float(ys_full.max()) - bf
            try:
                from scipy.signal import find_peaks
                pk, _pp = find_peaks(ys_full, height=bf + 0.10 * af,
                                     prominence=0.05 * af) if af > 0 else ([], None)
                npk = int(len(pk))
            except Exception:
                npk = 1
            # keep the full (sorted, smoothed) scan for the context overlay
            m_all, y_all = m.copy(), ys_full.copy()
            # optional analysis window — isolate ONE clean edge on a multi-peak scan.
            # The fields are ABSOLUTE axis positions (e.g. 27.433 and 27.438), not a
            # width; windowing is active only when max > min.
            lo, hi = float(win_lo.value), float(win_hi.value)
            win_intended = (lo > 0) or (hi > 0)
            windowed = hi > lo
            win_bad = win_intended and not windowed
            if windowed:
                sel = (m >= lo) & (m <= hi)
                if int(sel.sum()) < 5:
                    status.object = ("❌ Window selects <5 points — the fields are "
                                     "absolute axis values (e.g. 27.433 and 27.438), "
                                     "not a width. Check the bounds.")
                    return
                m, y = m[sel], y[sel]
            ys = gaussian_filter1d(y, float(smooth.value)) if smooth.value else y.copy()
            base = float(np.percentile(ys, 5))
            peak = float(ys.max())
            amp = peak - base
            if amp <= 0:
                status.object = "❌ Flat region — no edge to analyse."
                return
            ic = int(np.argmax(ys))
            u = ax  # axis unit name, e.g. "amu"; report in milli-units (m{u})
            warns = []
            if win_bad:
                warns.append(
                    f"⚠️ **Window ignored — full scan used.** *Window min* and "
                    f"*max* are absolute {u} positions and need max > min (e.g. "
                    f"27.433 and 27.438 to bracket the rising edge). You have "
                    f"min={lo:g}, max={hi:g}.")
            if npk > 1 and not windowed:
                warns.append(
                    f"⚠️ **{npk} peaks detected and no valid window.** The edge-"
                    f"finder splits at the tallest peak, so the ESF/LSF (and the "
                    f"slit-width number) are distorted by the interfering "
                    f"neighbouring peak. Window "
                    f"one *outer* edge — the lowest peak's rising edge or the "
                    f"highest peak's falling edge.")

            sides = (["rising", "falling"] if edge_sel.value == "both"
                     else [edge_sel.value])
            store = {}
            rows = []
            for side in sides:
                me, esf, lsf = _edge(m, ys, side, base, amp, ic)
                mom = _moments(me, lsf)
                try:
                    sigma, _info = fit_sigma_from_edge(m, y, side=side)
                    sigma = float(sigma)
                except Exception:
                    sigma = float("nan")
                store[side] = dict(me=me, esf=esf, lsf=lsf, mom=mom, sigma=sigma,
                                   m50=_interp_cross(me, esf, 0.5))
                fwhm = mom["fwhm"] * 1000
                sg = sigma * 1000 if np.isfinite(sigma) else float("nan")
                gfw = 2.3548 * sg if np.isfinite(sg) else float("nan")
                rows.append(
                    f"- **{side} edge** — LSF FWHM {fwhm:.3f} m{u} · "
                    f"skew {mom['skew']:+.2f} · excess kurtosis {mom['kurt']:+.2f} · "
                    f"fitted optics σ {sg:.3f} m{u} (Gaussian-equiv. FWHM {gfw:.3f} m{u})")

            # interference check: a clean single edge has a unimodal LSF. More than
            # one bump means an interfering neighbouring peak/shoulder is in range.
            pchk = "rising" if "rising" in store else sides[0]
            lc = store[pchk]["lsf"]
            try:
                from scipy.signal import find_peaks as _fp
                lpk, _ = _fp(lc, height=0.25 * (float(lc.max()) or 1.0),
                             distance=max(2, len(lc) // 50))
                if len(lpk) > 1:
                    warns.append(
                        f"⚠️ **The {pchk} edge's LSF has {len(lpk)} bumps** — it is "
                        f"not a single clean edge (a neighbouring peak or shoulder is "
                        f"inside the analysed region). Tighten the window to bracket "
                        f"only the outermost ramp; you'll know it's clean when the "
                        f"LSF is one bump and FWHM ≈ 2.355·σ.")
            except Exception:
                pass

            # consistency check: a clean Gaussian edge has FWHM ≈ 2.355·σ. A large
            # gap means a non-Gaussian / distorted edge even when the LSF is a single
            # bump — most often an INNER edge rising out of a valley (its foot sits on
            # a neighbour's tail) rather than an outer edge from true baseline.
            _mfw = store[pchk]["mom"]["fwhm"]
            _sg = store[pchk]["sigma"]
            if np.isfinite(_mfw) and np.isfinite(_sg) and _sg > 0:
                ratio = _mfw / (2.3548 * _sg)
                if ratio > 1.4 or ratio < 0.7:
                    warns.append(
                        f"⚠️ **LSF FWHM and 2.355·σ disagree by "
                        f"{abs(ratio - 1) * 100:.0f} %** ({_mfw * 1000:.3f} vs "
                        f"{2.3548 * _sg * 1000:.3f} m{u}) — the {pchk} edge isn't a "
                        f"clean single Gaussian. Most often this is an *inner* edge "
                        f"rising out of a valley (its foot sits on a neighbour's "
                        f"tail), not an *outer* edge from true baseline. Try the "
                        f"lowest peak's rising edge or the highest peak's falling "
                        f"edge; a genuine gap that survives an outer edge is real "
                        f"aberration.")
            peaks_warn.object = "\n\n".join(warns)

            # slit / flat-top width = separation of one peak's two edge midpoints.
            # Only valid for a single isolated peak (or a window around one peak),
            # and only when both edges were extracted.
            slit_w = None
            if len(store) == 2 and (npk == 1 or windowed):
                m50_r = _interp_cross(*(_edge(m, ys, "rising", base, amp, ic)[:2]), 0.5)
                m50_f = _interp_cross(*(_edge(m, ys, "falling", base, amp, ic)[:2]), 0.5)
                slit_w = abs(m50_f - m50_r) * 1000

            # optics-limited MRP from the analysed edge — same 10–90% definition
            # and M as the HMR tab, so it's directly comparable on the same edge.
            pmrp = "rising" if "rising" in store else sides[0]
            _me, _esf = store[pmrp]["me"], store[pmrp]["esf"]
            _m10 = _interp_cross(_me, _esf, 0.10)
            _m90 = _interp_cross(_me, _esf, 0.90)
            dM_clean = abs(_m90 - _m10)
            M_nom = int(round(float(np.median(m_all))))
            mrp_clean = (M_nom / dM_clean) if dM_clean > 0 else float("nan")

            # abundance sensitivity: tail level at the analysed peak ± Δm, ppm of peak
            center = float(m[ic])
            dm = float(dm_as.value) / 1000.0
            ybl = np.clip(ys - base, 0, None)
            as_lo = float(np.interp(center - dm, m, ybl)) / amp * 1e6
            as_hi = float(np.interp(center + dm, m, ybl)) / amp * 1e6

            # ---- figure 0: context — full scan with the analysed region shaded --
            figs = []
            figc, axc = plt.subplots(figsize=(6.2, 2.6))
            axc.plot(m_all, np.clip(y_all, 0.5, None), color="#888", lw=1.0)
            axc.axvspan(float(m.min()), float(m.max()), color="#1f77b4",
                        alpha=0.15, label="analysed region")
            axc.set_yscale("log")
            axc.set_xlabel(f"mass ({u}, nominal)")
            axc.set_ylabel("counts")
            axc.set_title("Scan — analysed region shaded "
                          + ("(windowed)" if windowed else "(FULL scan)"),
                          fontsize=9)
            axc.legend(fontsize=8, loc="upper right")
            figc.tight_layout()
            figs.append(("context", _fig_png(figc)))

            # ---- figure 1: ESF + LSF for each side, centred on its midpoint ----
            fig, ax1 = plt.subplots(figsize=(6.2, 4.0))
            ax2 = ax1.twinx()
            colors = {"rising": "#1f77b4", "falling": "#d62728"}
            for side, d in store.items():
                x = (d["me"] - d["m50"]) * 1000
                ax1.plot(x, d["esf"], color=colors[side], lw=1.6,
                         label=f"ESF ({side})")
                lsfn = d["lsf"] / (d["lsf"].max() or 1.0)
                ax2.plot(x, lsfn, color=colors[side], lw=1.2, ls="--", alpha=0.8,
                         label=f"LSF ({side})")
            ax1.set_xlabel(f"Δm from edge midpoint (m{u})")
            ax1.set_ylabel("edge-spread (norm.)")
            ax2.set_ylabel("line-spread (norm.)")
            ax1.set_title(f"ESF & LSF — {HMR_SHARED.get('name')}", fontsize=10)
            ax1.grid(alpha=0.25)
            h1, l1 = ax1.get_legend_handles_labels()
            h2, l2 = ax2.get_legend_handles_labels()
            ax1.legend(h1 + h2, l1 + l2, fontsize=8, loc="upper right")
            fig.tight_layout()
            figs.append(("ESF & LSF", _fig_png(fig)))

            # ---- figure 2: rising/falling LSF overlay (asymmetry) if both ----
            if len(store) == 2:
                fig2, axo = plt.subplots(figsize=(6.2, 3.6))
                for side, d in store.items():
                    x = (d["me"] - d["m50"]) * 1000
                    if side == "falling":
                        x = -x   # mirror so both edges rise in the same sense
                    lsfn = d["lsf"] / (d["lsf"].max() or 1.0)
                    axo.plot(x, lsfn, color=colors[side], lw=1.5, label=side)
                axo.set_xlabel(f"|Δm| from edge midpoint (m{u})")
                axo.set_ylabel("LSF (norm.)")
                axo.set_title("Rising vs falling LSF — gap = aberration/asymmetry",
                              fontsize=10)
                axo.grid(alpha=0.25); axo.legend(fontsize=8)
                fig2.tight_layout()
                figs.append(("Rising/falling overlay", _fig_png(fig2)))

            # ---- figure 3: MTF = |FFT(LSF)| (rising or first side) ----
            pside = "rising" if "rising" in store else sides[0]
            d = store[pside]
            me = d["me"]
            # resample LSF onto a uniform grid before FFT
            n = max(64, len(me))
            mu_grid = np.linspace(me.min(), me.max(), n)
            lsf_u = np.interp(mu_grid, me, d["lsf"])
            lsf_u = lsf_u - lsf_u.min()
            dmu = (mu_grid[1] - mu_grid[0])
            mtf = np.abs(np.fft.rfft(lsf_u))
            mtf = mtf / (mtf[0] or 1.0)
            freq = np.fft.rfftfreq(n, d=dmu)   # cycles per axis-unit
            fig3, axm = plt.subplots(figsize=(6.2, 3.4))
            axm.plot(freq, mtf, color="#2ca02c", lw=1.6)
            axm.set_xlabel(f"spatial frequency (cycles per {u})")
            axm.set_ylabel("MTF (norm.)")
            axm.set_title(f"MTF (FFT of {pside} LSF)", fontsize=10)
            axm.set_ylim(0, 1.02); axm.grid(alpha=0.25)
            fig3.tight_layout()
            figs.append(("MTF", _fig_png(fig3)))

            fig_box.objects = [
                pn.pane.PNG(p, width=560) for _, p in figs]

            slit_line = (
                f"\n\n**Slit / flat-top width** (edge-midpoint separation) = "
                f"{slit_w:.3f} m{u} — the slit's contribution; the LSF FWHM above "
                f"is the optics' contribution."
                if slit_w is not None else
                "\n\n*Slit / flat-top width not shown — it's only meaningful for a "
                "single isolated peak with both edges (use Edge = both on one "
                "windowed peak).*")
            metrics.object = (
                "**LSF metrics**\n\n"
                + "\n".join(rows)
                + f"\n\n**Optics-limited MRP** (10–90 % of the {pmrp} edge) = "
                f"**{mrp_clean:,.0f}** (M = {M_nom}, ΔM = {dM_clean * 1000:.3f} "
                f"m{u}) — set by the beam, independent of slit width. Same "
                f"definition as the HMR tab's MRP, so compare them directly on the "
                f"same edge; a much lower HMR value means its edge includes an "
                f"interfering mass peak."
                + slit_line
                + f"\n\n**Abundance sensitivity** at Δm = {dm_as.value:.1f} m{u} "
                f"from the analysed peak: low-mass side {as_lo:,.0f} ppm, "
                f"high-mass side {as_hi:,.0f} ppm (tail counts as a fraction of "
                f"peak height — the skirt a minor isotope would sit on; on a real "
                f"multi-peak scan the measured valley between peaks is the more "
                f"direct number).")
            status.object = f"✅ Computed LSF for **{HMR_SHARED.get('name')}**."
        except Exception as exc:
            status.object = f"❌ LSF computation failed: `{exc}`"
    compute_btn.on_click(_compute)

    note = pn.pane.Markdown(
        "<span style='font-size:11px;opacity:0.7'>The LSF FWHM and the fitted optics "
        "σ measure the same thing two ways: for a clean Gaussian beam, LSF FWHM ≈ "
        "2.355·σ. A large gap between them, non-zero skew, or rising/falling edges "
        "that don't overlap all point to aberration rather than a simple defocus. "
        "Abundance sensitivity is the practical pay-off — the tail level set here is "
        "the floor a neighbouring minor isotope (e.g. ¹²C¹⁵N next to ¹³C¹⁴N) has to "
        "be seen above. For a multi-peak scan (e.g. mass-27 ¹²C¹⁵N next to "
        "¹³C¹⁴N), set the **window** around a single *outer* edge — the clean "
        "single-species edges are the lowest peak's rising edge and the highest "
        "peak's falling edge. Series mode (metric vs tuning setting) is the "
        "planned next step.</span>")

    return pn.Column(
        intro, pn.Row(src_btn), status, pn.layout.Divider(),
        pn.Row(axis_sel, edge_sel, smooth),
        pn.Row(win_lo, win_hi, dm_as),
        pn.Row(compute_btn),
        peaks_warn, metrics, fig_box, note, sizing_mode="stretch_width")


# ─────────────────────────────────────────────────────────────────────────────
# Assemble the app
# ─────────────────────────────────────────────────────────────────────────────

tabs = pn.Tabs(
    ("Imaging",  pn.panel(pn.bind(imaging_view, state.param.img))),
    ("HMR",      hmr_view()),
    ("Ion optics", lsf_view()),
    ("Analysis", pn.panel(pn.bind(analysis_view, state.param.img))),
    ("Draw ROIs", pn.panel(pn.bind(draw_roi_view, state.param.img))),
    ("3D",       pn.panel(pn.bind(three_d_view, state.param.img))),
    ("Metadata", pn.panel(pn.bind(metadata_view, state.param.img))),
    ("Help",     help_view()),
    dynamic=True,
    sizing_mode="stretch_width",
)

sidebar = pn.Column(
    pn.pane.Markdown("### Load a file"),
    dir_input,
    up_btn,
    subdir_select,
    file_select,
    load_btn,
    load_status,
    pn.layout.Divider(),
    bin_select,
    bin_note,
    pn.layout.Divider(),
    pn.pane.Markdown("### Session"),
    session_dir,
    session_name,
    save_session_btn,
    session_select,
    load_session_btn,
    session_status,
    session_note,
    pn.layout.Divider(),
    pn.pane.Markdown("### Journal"),
    journal_status,
    journal_dir,
    journal_name,
    journal_fmt,
    export_journal_btn,
    pn.Row(remove_last_btn, clear_journal_btn),
    export_status,
    journal_note,
    pn.layout.Divider(),
    pn.pane.Markdown(
        f"*PyMIMS v{PYMIMS_VERSION} — NanoSIMS .im toolkit*  \n"
        "<span style='font-size:11px;opacity:0.7'>Free &amp; open-source. If it "
        "saves you time, you can support development at "
        "<a href='https://n1lab.substack.com' target='_blank'>The N=1 Lab</a>."
        "</span>"),
)

template = pn.template.FastListTemplate(
    title="PyMIMS",
    theme="dark",
    header_background="#121212",
    sidebar=[sidebar],
    main=[tabs],
    sidebar_width=340,
)

template.servable()


if __name__ == "__main__":
    pn.serve(template, show=True, port=5006)
